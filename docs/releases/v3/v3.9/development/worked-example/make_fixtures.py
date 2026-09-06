#!/usr/bin/env python3
"""Generate the two sample source documents for the /presentify worked example.

Writes a small PowerPoint deck (single-deck "preserve the flow" mode) and a
Word report (single-report "present the report" mode) into ./inputs/. These
are demo INPUTS, not distributed catalog content, so they live under
docs/v3/v3.9/development/ and are git-ignored (regenerate with this script).

Local-only; uses python-pptx, python-docx, and Pillow. No network calls.
"""
from __future__ import annotations

from pathlib import Path

HERE = Path(__file__).resolve().parent
INPUTS = HERE / "inputs"


def make_logo(path: Path) -> None:
    """Draw a tiny, intentional placeholder mark so the image block is exercised."""
    from PIL import Image, ImageDraw

    img = Image.new("RGB", (160, 120), (16, 24, 40))
    draw = ImageDraw.Draw(img)
    # A small ascending three-bar motif (a "growth" glyph), not a logo of any product.
    bars = [(28, 86, 56, 100), (66, 60, 94, 100), (104, 34, 132, 100)]
    for i, (x0, y0, x1, y1) in enumerate(bars):
        shade = (90 + i * 50, 150 + i * 30, 220)
        draw.rectangle([x0, y0, x1, y1], fill=shade)
    img.save(path, "PNG")


def make_pptx(path: Path, logo: Path) -> None:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    blank_title = prs.slide_layouts[0]
    title_content = prs.slide_layouts[1]

    # Slide 1 -- title
    s = prs.slides.add_slide(blank_title)
    s.shapes.title.text = "Quarterly Business Review"
    s.placeholders[1].text = "FY26 Q3 - Operations and Growth"

    # Slide 2 -- agenda (bullets)
    s = prs.slides.add_slide(title_content)
    s.shapes.title.text = "Agenda"
    tf = s.placeholders[1].text_frame
    tf.text = "Where we are this quarter"
    for line in ["Revenue by region", "Key operational highlights", "Next steps"]:
        p = tf.add_paragraph()
        p.text = line
    s.notes_slide.notes_text_frame.text = (
        "Keep the open tight: one sentence per agenda item, then move on."
    )

    # Slide 3 -- data table (the enrichment pass promotes this to a chart) + image
    s = prs.slides.add_slide(title_content)
    s.shapes.title.text = "Revenue by Region"
    rows, cols = 4, 3
    table = s.shapes.add_table(
        rows, cols, Inches(0.6), Inches(1.6), Inches(5.2), Inches(2.4)
    ).table
    headers = ["Region", "Q2 ($M)", "Q3 ($M)"]
    for c, h in enumerate(headers):
        table.cell(0, c).text = h
    data = [("North America", "120", "145"), ("EMEA", "85", "98"), ("APAC", "60", "77")]
    for r, (region, q2, q3) in enumerate(data, start=1):
        table.cell(r, 0).text = region
        table.cell(r, 1).text = q2
        table.cell(r, 2).text = q3
    s.shapes.add_picture(str(logo), Inches(6.2), Inches(1.8), height=Inches(1.4))

    # Slide 4 -- nested highlights + notes
    s = prs.slides.add_slide(title_content)
    s.shapes.title.text = "Key Highlights"
    tf = s.placeholders[1].text_frame
    tf.text = "Growth held across every region"
    nested = [
        ("APAC outpaced plan", 0),
        ("Driven by two enterprise wins", 1),
        ("Margin improved 3 points", 0),
        ("Lower cloud spend per request", 1),
        ("Churn fell to a new low", 0),
    ]
    for text, level in nested:
        p = tf.add_paragraph()
        p.text = text
        p.level = level
        for run in p.runs:
            run.font.size = Pt(18)
    s.notes_slide.notes_text_frame.text = (
        "Land the margin point: it is the headline for the board."
    )

    # Slide 5 -- next steps
    s = prs.slides.add_slide(title_content)
    s.shapes.title.text = "Next Steps"
    tf = s.placeholders[1].text_frame
    tf.text = "Three commitments for Q4"
    for line in [
        "Expand the APAC enterprise motion",
        "Hold the margin gains while scaling",
        "Ship the reliability roadmap",
    ]:
        p = tf.add_paragraph()
        p.text = line

    prs.save(str(path))


def make_docx(path: Path) -> None:
    from docx import Document

    doc = Document()
    doc.add_heading("Operational Readiness Report", level=0)

    doc.add_heading("Executive Summary", level=1)
    doc.add_paragraph(
        "This report assesses operational readiness for the next release cycle. "
        "Reliability and incident response both improved against target, and the "
        "remaining risks are tracked with named owners."
    )

    doc.add_heading("Findings", level=1)
    doc.add_paragraph(
        "Service reliability exceeded its target for the quarter, and incident "
        "volume fell in every quarter of the year."
    )
    for item in [
        "Uptime cleared the 99.9 percent target with room to spare",
        "Recovery time dropped after the on-call runbook refresh",
        "Incident volume is at its lowest in four quarters",
    ]:
        doc.add_paragraph(item, style="List Bullet")
    table = doc.add_table(rows=5, cols=2)
    table.style = "Light Grid Accent 1"
    hdr = table.rows[0].cells
    hdr[0].text, hdr[1].text = "Quarter", "Incidents"
    trend = [("Q1", "18"), ("Q2", "14"), ("Q3", "12"), ("Q4", "7")]
    for r, (quarter, incidents) in enumerate(trend, start=1):
        cells = table.rows[r].cells
        cells[0].text, cells[1].text = quarter, incidents

    doc.add_heading("Recommendations", level=1)
    for item in [
        "Make the runbook refresh a standing quarterly task",
        "Add a synthetic check for the slowest recovery path",
        "Rehearse a regional failover before the next release",
    ]:
        doc.add_paragraph(item, style="List Bullet")

    doc.add_heading("Conclusion", level=1)
    doc.add_paragraph(
        "Readiness is on track. Sustaining the runbook discipline is the single "
        "highest-leverage action going into the next cycle."
    )

    doc.save(str(path))


def main() -> None:
    INPUTS.mkdir(parents=True, exist_ok=True)
    logo = INPUTS / "growth-mark.png"
    make_logo(logo)
    make_pptx(INPUTS / "sample-deck.pptx", logo)
    make_docx(INPUTS / "sample-report.docx")
    print(f"Wrote fixtures to {INPUTS}")
    for p in sorted(INPUTS.iterdir()):
        print(f"  {p.name} ({p.stat().st_size} bytes)")


if __name__ == "__main__":
    main()
