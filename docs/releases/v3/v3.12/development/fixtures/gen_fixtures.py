#!/usr/bin/env python3
"""Generate the Phase 1 extraction-fidelity fixtures into ./inputs/.

Fixtures (all synthetic, deterministic content):
    deck.pdf    - landscape PDF-from-slides proxy: vector bar chart, vector
                  map, embedded photo, repeated logo (5 pages), captions.
    deck.pptx   - grouped shape, native bar chart, picture, table, notes.
    report.docx - Title/Heading structure, list, table, inline image, and an
                  injected native chart part (word/charts/chart1.xml).
    data.xlsx   - one numeric range (chart-block source).
    scanned.pdf - image-only pages (no text layer) with known text, a table,
                  and a bar figure, for the two-tier OCR path.

Ground-truth values asserted by verify_phase1.py:
    deck.pdf chart bars: Q1 120, Q2 135, Q3 150, Q4 170
    deck.pptx chart:     Q1 120, Q2 135, Q3 150, Q4 170 (series "Revenue")
    report.docx chart:   North 10, South 20, East 30 (series "Units")
    data.xlsx range:     North 42, South 37, East 55 (series "Revenue")
    scanned.pdf text:    "QUARTERLY UPDATE", "Revenue grew twelve percent",
                         table North 42 / South 37.
"""

from __future__ import annotations

import io
import zipfile
from pathlib import Path

HERE = Path(__file__).resolve().parent
INPUTS = HERE / "inputs"

DECK_CATEGORIES = ["Q1", "Q2", "Q3", "Q4"]
DECK_VALUES = [120.0, 135.0, 150.0, 170.0]
DOCX_CATEGORIES = ["North", "South", "East"]
DOCX_VALUES = [10.0, 20.0, 30.0]
XLSX_ROWS = [("North", 42), ("South", 37), ("East", 55)]


def _photo_png(width: int = 400, height: int = 267) -> bytes:
    """Deterministic gradient 'photo' PNG."""
    from PIL import Image

    image = Image.new("RGB", (width, height))
    pixels = image.load()
    for y in range(height):
        for x in range(width):
            pixels[x, y] = (40 + x * 180 // width, 80 + y * 120 // height, 120)
    buffer = io.BytesIO()
    image.save(buffer, format="PNG")
    return buffer.getvalue()


def _logo_png() -> bytes:
    """Small deterministic logo PNG (repeated on every deck.pdf page)."""
    from PIL import Image, ImageDraw

    image = Image.new("RGB", (120, 60), (20, 60, 130))
    draw = ImageDraw.Draw(image)
    draw.rectangle([8, 8, 112, 52], outline=(255, 255, 255), width=4)
    draw.ellipse([44, 16, 76, 44], fill=(255, 200, 40))
    buffer = io.BytesIO()
    image.save(buffer, format="PNG")
    return buffer.getvalue()


def make_deck_pdf(path: Path) -> None:
    from reportlab.lib.pagesizes import landscape, letter
    from reportlab.lib.utils import ImageReader
    from reportlab.pdfgen import canvas

    page_size = landscape(letter)  # 792 x 612 pt
    width, height = page_size
    logo = ImageReader(io.BytesIO(_logo_png()))
    photo = ImageReader(io.BytesIO(_photo_png()))
    pdf = canvas.Canvas(str(path), pagesize=page_size)

    def chrome(title: str) -> None:
        pdf.drawImage(logo, width - 100, height - 50, width=60, height=30)
        pdf.setFont("Helvetica-Bold", 30)
        pdf.drawString(50, height - 70, title)
        pdf.setFont("Helvetica", 12)

    # Page 1 - title slide.
    chrome("Nexus Board Review")
    pdf.setFont("Helvetica", 16)
    pdf.drawString(50, height - 110, "Fiscal year 2026 highlights")
    pdf.showPage()

    # Page 2 - vector bar chart (rects + axis lines) with a caption.
    chrome("Revenue by Quarter")
    base_y, bar_width, gap, x_start = 170.0, 60.0, 30.0, 130.0
    pdf.setLineWidth(1.5)
    pdf.line(110, base_y, 520, base_y)  # x axis
    pdf.line(110, base_y, 110, 440)  # y axis
    pdf.setFillColorRGB(0.15, 0.35, 0.65)
    for index, value in enumerate(DECK_VALUES):
        x = x_start + index * (bar_width + gap)
        pdf.rect(x, base_y, bar_width, value * 1.4, stroke=0, fill=1)
    pdf.setFillColorRGB(0, 0, 0)
    pdf.setFont("Helvetica", 10)
    for index, label in enumerate(DECK_CATEGORIES):
        x = x_start + index * (bar_width + gap) + bar_width / 2 - 8
        pdf.drawString(x, base_y - 14, label)
    for tick in (0, 50, 100, 150):
        pdf.drawString(84, base_y + tick * 1.4 - 3, str(tick))
    pdf.setFont("Helvetica", 12)
    pdf.drawString(110, 132, "Figure 1: Revenue by quarter (USD millions)")
    pdf.showPage()

    # Page 3 - vector 'map' (bezier outline + site dots) with a caption.
    chrome("Enrollment Map")
    pdf.setLineWidth(2)
    pdf.bezier(160, 200, 240, 420, 380, 440, 470, 380)
    pdf.bezier(470, 380, 560, 320, 520, 230, 430, 210)
    pdf.bezier(430, 210, 340, 180, 220, 170, 160, 200)
    pdf.setFillColorRGB(0.75, 0.2, 0.2)
    for x, y in ((250, 300), (330, 340), (410, 280), (470, 330)):
        pdf.circle(x, y, 6, stroke=0, fill=1)
    pdf.setFillColorRGB(0, 0, 0)
    pdf.drawString(160, 138, "Figure 2: Site enrollment map")
    pdf.showPage()

    # Page 4 - embedded photo with a caption.
    chrome("Team")
    pdf.drawImage(photo, 150, 200, width=300, height=200)
    pdf.drawString(150, 172, "Photo: Team offsite, June 2026")
    pdf.showPage()

    # Page 5 - summary text.
    chrome("Summary")
    for offset, line in enumerate(
        (
            "Revenue reached 170 in Q4.",
            "Enrollment expanded to four sites.",
            "Next review scheduled for January.",
        )
    ):
        pdf.drawString(70, height - 130 - offset * 22, line)
    pdf.showPage()
    pdf.save()


def make_deck_pptx(path: Path) -> None:
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.util import Inches

    presentation = Presentation()
    blank = presentation.slide_layouts[6]
    title_layout = presentation.slide_layouts[0]

    slide = presentation.slides.add_slide(title_layout)
    slide.shapes.title.text = "Nexus Board Review"
    slide.placeholders[1].text = "Fiscal year 2026 highlights"

    slide = presentation.slides.add_slide(blank)
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(6), Inches(0.8))
    box.text_frame.text = "Agenda"
    body = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(6), Inches(2))
    body.text_frame.text = "Revenue"
    body.text_frame.add_paragraph().text = "Enrollment"
    body.text_frame.add_paragraph().text = "Team"

    # Grouped shape holding a textbox (group-recursion check).
    slide = presentation.slides.add_slide(blank)
    head = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(6), Inches(0.8))
    head.text_frame.text = "Key Insight"
    group = slide.shapes.add_group_shape()
    grouped_box = group.shapes.add_textbox(
        Inches(1.0), Inches(2.0), Inches(6.0), Inches(1.0)
    )
    grouped_box.text_frame.text = "Grouped insight: retention doubled year over year"

    # Native chart (real series values).
    slide = presentation.slides.add_slide(blank)
    head = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(6), Inches(0.8))
    head.text_frame.text = "Revenue by Quarter"
    chart_data = CategoryChartData()
    chart_data.categories = DECK_CATEGORIES
    chart_data.add_series("Revenue", tuple(DECK_VALUES))
    frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1),
        Inches(1.5),
        Inches(7),
        Inches(4.5),
        chart_data,
    )
    chart = frame.chart
    chart.has_title = True
    chart.chart_title.text_frame.text = "Revenue by quarter (USD millions)"

    # Picture + table + notes.
    slide = presentation.slides.add_slide(blank)
    head = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(6), Inches(0.8))
    head.text_frame.text = "Team and Sites"
    photo_path = path.parent / "_photo_tmp.png"
    photo_path.write_bytes(_photo_png())
    slide.shapes.add_picture(str(photo_path), Inches(0.7), Inches(1.4), Inches(3))
    table = slide.shapes.add_table(
        3, 2, Inches(4.5), Inches(1.4), Inches(4), Inches(1.5)
    ).table
    table.cell(0, 0).text = "Site"
    table.cell(0, 1).text = "Enrolled"
    table.cell(1, 0).text = "North"
    table.cell(1, 1).text = "42"
    table.cell(2, 0).text = "South"
    table.cell(2, 1).text = "37"
    slide.notes_slide.notes_text_frame.text = "Speaker note: thank the site leads."
    head.text_frame.word_wrap = True

    presentation.save(str(path))
    photo_path.unlink(missing_ok=True)


_CHART_XML = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<c:chartSpace xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <c:chart>
    <c:plotArea>
      <c:layout/>
      <c:barChart>
        <c:barDir val="col"/>
        <c:grouping val="clustered"/>
        <c:ser>
          <c:idx val="0"/>
          <c:order val="0"/>
          <c:tx><c:strRef><c:f>Sheet1!$B$1</c:f><c:strCache><c:ptCount val="1"/><c:pt idx="0"><c:v>Units</c:v></c:pt></c:strCache></c:strRef></c:tx>
          <c:cat><c:strRef><c:f>Sheet1!$A$2:$A$4</c:f><c:strCache><c:ptCount val="3"/><c:pt idx="0"><c:v>North</c:v></c:pt><c:pt idx="1"><c:v>South</c:v></c:pt><c:pt idx="2"><c:v>East</c:v></c:pt></c:strCache></c:strRef></c:cat>
          <c:val><c:numRef><c:f>Sheet1!$B$2:$B$4</c:f><c:numCache><c:formatCode>General</c:formatCode><c:ptCount val="3"/><c:pt idx="0"><c:v>10</c:v></c:pt><c:pt idx="1"><c:v>20</c:v></c:pt><c:pt idx="2"><c:v>30</c:v></c:pt></c:numCache></c:numRef></c:val>
        </c:ser>
      </c:barChart>
    </c:plotArea>
  </c:chart>
</c:chartSpace>
"""

_CHART_DRAWING = (
    '<w:p xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
    "<w:r><w:drawing>"
    '<wp:inline distT="0" distB="0" distL="0" distR="0" '
    'xmlns:wp="http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing">'
    '<wp:extent cx="4572000" cy="2743200"/>'
    '<wp:docPr id="99" name="Chart 99"/>'
    '<a:graphic xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
    '<a:graphicData uri="http://schemas.openxmlformats.org/drawingml/2006/chart">'
    '<c:chart xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart" '
    'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" '
    'r:id="rId100"/>'
    "</a:graphicData></a:graphic></wp:inline></w:drawing></w:r></w:p>"
)


def make_report_docx(path: Path) -> None:
    from docx import Document
    from docx.shared import Inches

    document = Document()
    document.add_heading("Nexus Annual Report", level=0)  # Title style
    document.add_heading("Introduction", level=1)
    document.add_paragraph(
        "This report summarizes the fiscal year across revenue, enrollment, "
        "and team growth."
    )
    document.add_heading("Results", level=1)
    document.add_paragraph("Revenue grew in every region.", style="List Bullet")
    document.add_paragraph("Enrollment doubled at two sites.", style="List Bullet")
    table = document.add_table(rows=3, cols=2)
    table.rows[0].cells[0].text = "Region"
    table.rows[0].cells[1].text = "Revenue"
    table.rows[1].cells[0].text = "North"
    table.rows[1].cells[1].text = "42"
    table.rows[2].cells[0].text = "South"
    table.rows[2].cells[1].text = "37"
    photo_path = path.parent / "_photo_tmp.png"
    photo_path.write_bytes(_photo_png())
    document.add_picture(str(photo_path), width=Inches(3))
    document.save(str(path))
    photo_path.unlink(missing_ok=True)

    # Inject a native chart part (python-docx cannot author charts).
    raw = path.read_bytes()
    source = zipfile.ZipFile(io.BytesIO(raw))
    buffer = io.BytesIO()
    with zipfile.ZipFile(buffer, "w", zipfile.ZIP_DEFLATED) as target:
        for item in source.infolist():
            data = source.read(item.filename)
            if item.filename == "[Content_Types].xml":
                text = data.decode("utf-8").replace(
                    "</Types>",
                    '<Override PartName="/word/charts/chart1.xml" ContentType='
                    '"application/vnd.openxmlformats-officedocument.drawingml.'
                    'chart+xml"/></Types>',
                )
                data = text.encode("utf-8")
            elif item.filename == "word/_rels/document.xml.rels":
                text = data.decode("utf-8").replace(
                    "</Relationships>",
                    '<Relationship Id="rId100" Type="http://schemas.'
                    "openxmlformats.org/officeDocument/2006/relationships/chart"
                    '" Target="charts/chart1.xml"/></Relationships>',
                )
                data = text.encode("utf-8")
            elif item.filename == "word/document.xml":
                text = data.decode("utf-8")
                marker = "<w:sectPr"
                position = text.find(marker)
                if position == -1:
                    position = text.find("</w:body>")
                data = (text[:position] + _CHART_DRAWING + text[position:]).encode(
                    "utf-8"
                )
            target.writestr(item, data)
        target.writestr("word/charts/chart1.xml", _CHART_XML)
    source.close()
    path.write_bytes(buffer.getvalue())


def make_data_xlsx(path: Path) -> None:
    from openpyxl import Workbook

    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Revenue"
    sheet.append(["Region", "Revenue"])
    for region, value in XLSX_ROWS:
        sheet.append([region, value])
    workbook.save(str(path))


def make_scanned_pdf(path: Path) -> None:
    from PIL import Image, ImageDraw, ImageFont

    def _font(size: int):
        candidates = (
            "arial.ttf",  # Windows
            "DejaVuSans.ttf",
            "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",  # Ubuntu CI
            "/System/Library/Fonts/Supplemental/Arial.ttf",  # macOS
        )
        for name in candidates:
            try:
                return ImageFont.truetype(name, size)
            except OSError:
                continue
        return ImageFont.load_default()

    heading_font = _font(48)
    body_font = _font(30)

    def page_one() -> Image.Image:
        image = Image.new("RGB", (1700, 1300), "white")
        draw = ImageDraw.Draw(image)
        draw.text((100, 80), "QUARTERLY UPDATE", font=heading_font, fill="black")
        draw.text(
            (100, 200),
            "Revenue grew twelve percent over the prior quarter.",
            font=body_font,
            fill="black",
        )
        draw.text(
            (100, 260),
            "Enrollment reached both regional targets.",
            font=body_font,
            fill="black",
        )
        # Simple two-column table with a grid.
        rows = [("Region", "Units"), ("North", "42"), ("South", "37")]
        top, left, row_h, col_w = 400, 100, 70, 320
        for row_index, (left_cell, right_cell) in enumerate(rows):
            y = top + row_index * row_h
            draw.rectangle([left, y, left + 2 * col_w, y + row_h], outline="black")
            draw.line([left + col_w, y, left + col_w, y + row_h], fill="black")
            draw.text((left + 20, y + 15), left_cell, font=body_font, fill="black")
            draw.text(
                (left + col_w + 20, y + 15), right_cell, font=body_font, fill="black"
            )
        return image

    def page_two() -> Image.Image:
        image = Image.new("RGB", (1700, 1300), "white")
        draw = ImageDraw.Draw(image)
        draw.text((100, 80), "REVENUE FIGURE", font=heading_font, fill="black")
        base, bar_w, gap = 1000, 160, 90
        for index, value in enumerate((300, 420, 520)):
            x = 200 + index * (bar_w + gap)
            draw.rectangle([x, base - value, x + bar_w, base], fill=(30, 80, 160))
        draw.text(
            (200, 1040),
            "Figure: revenue by region (illustrative)",
            font=body_font,
            fill="black",
        )
        return image

    first, second = page_one(), page_two()
    first.save(str(path), save_all=True, append_images=[second], resolution=150.0)


def main() -> int:
    INPUTS.mkdir(parents=True, exist_ok=True)
    make_deck_pdf(INPUTS / "deck.pdf")
    make_deck_pptx(INPUTS / "deck.pptx")
    make_report_docx(INPUTS / "report.docx")
    make_data_xlsx(INPUTS / "data.xlsx")
    make_scanned_pdf(INPUTS / "scanned.pdf")
    names = sorted(p.name for p in INPUTS.iterdir() if not p.name.startswith("_"))
    print(f"Generated {len(names)} fixtures: {', '.join(names)}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
