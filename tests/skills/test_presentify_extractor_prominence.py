"""Confirm the v3.15.4 Phase 2 extractor prominence signals.

The image-sizing discipline in `references/interactive-features.md` is
data-driven: it ranks visuals by `page_fraction` (share of the source
page/slide area) and native `width`/`height`. These tests confirm
`extract_content.py` populates all three for BOTH a PPTX picture shape and a PDF
embedded raster, so the ranking is not a guess.

The extractor and its fixture builders need heavy optional libraries
(python-pptx / reportlab / pdfplumber / Pillow) that the general `ci.yml` tests
job does not install; each test `importorskip`s the libraries it needs, so it
runs fully in the `presentify-extractor` workflow (which installs them) and
skips cleanly elsewhere. The extractor is loaded by path via importlib, matching
the pattern in test_media_key_setup.py.
"""

from __future__ import annotations

import importlib.util
import io
import json
from pathlib import Path

import pytest

pytest.importorskip("PIL", reason="Pillow is required to build image fixtures")

_ROOT = Path(__file__).resolve().parents[2]
_EXTRACT_PATH = (
    _ROOT
    / "catalog"
    / "skills"
    / "specialized-domains"
    / "document-to-interactive-html"
    / "scripts"
    / "extract_content.py"
)


def _load(path: Path, name: str):
    spec = importlib.util.spec_from_file_location(name, path)
    assert spec and spec.loader, f"cannot load {path}"
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)
    return module


extract = _load(_EXTRACT_PATH, "extract_content")


def _photo_png(width: int, height: int) -> bytes:
    """A deterministic RGB gradient PNG of the requested native size."""
    from PIL import Image

    image = Image.new("RGB", (width, height))
    pixels = image.load()
    for y in range(height):
        for x in range(width):
            pixels[x, y] = (
                40 + x * 180 // width,
                80 + y * 120 // height,
                120,
            )
    buffer = io.BytesIO()
    image.save(buffer, "PNG")
    return buffer.getvalue()


def _image_blocks(model: dict) -> list[dict]:
    return [
        block
        for section in model.get("sections", [])
        for block in section.get("blocks", [])
        if block.get("type") == "image"
    ]


def _extract(path: Path, tmp_path: Path) -> dict:
    out = tmp_path / "model.json"
    rc = extract.main([str(path), "-o", str(out)])
    assert rc == 0, f"extractor exited {rc} for {path.name}"
    return json.loads(out.read_text(encoding="utf-8"))


# --- PPTX picture shape -----------------------------------------------------


def test_pptx_picture_has_prominence_signals(tmp_path):
    pptx = pytest.importorskip("pptx")
    from pptx.util import Inches

    photo = tmp_path / "photo.png"
    photo.write_bytes(_photo_png(400, 267))
    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    slide.shapes.add_picture(str(photo), Inches(1), Inches(1), Inches(4))
    deck = tmp_path / "deck.pptx"
    prs.save(str(deck))

    images = _image_blocks(_extract(deck, tmp_path))
    assert images, "no image block extracted from the PPTX fixture"
    block = next((b for b in images if b.get("origin") == "shape-picture"), images[0])

    # Native dimensions preserved from the source bytes.
    assert block.get("width") == 400
    assert block.get("height") == 267
    # page_fraction is populated and a sane 0..1 share of the slide area.
    frac = block.get("page_fraction")
    assert isinstance(frac, (int, float)), "PPTX picture missing page_fraction"
    assert 0.0 < frac <= 1.0


# --- PDF embedded raster ----------------------------------------------------


def test_pdf_embedded_raster_has_prominence_signals(tmp_path):
    # A PDF with a REAL text layer plus an embedded raster exercises the
    # embedded-raster path (a text-free image-only page is classified
    # scanned-page, which correctly carries no page_fraction).
    pytest.importorskip("pdfplumber")
    pytest.importorskip("pypdf")
    reportlab = pytest.importorskip("reportlab")
    del reportlab
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.utils import ImageReader
    from reportlab.pdfgen import canvas

    pdf = tmp_path / "report.pdf"
    width, height = letter
    doc = canvas.Canvas(str(pdf), pagesize=letter)
    doc.setFont("Helvetica-Bold", 20)
    doc.drawString(50, height - 70, "Quarterly Report")
    doc.setFont("Helvetica", 12)
    for index, line in enumerate(
        ["Line one of body text.", "Line two with detail.", "Line three closing."]
    ):
        doc.drawString(50, height - 110 - index * 22, line)
    doc.drawImage(
        ImageReader(io.BytesIO(_photo_png(600, 400))), 150, 200, width=300, height=200
    )
    doc.showPage()
    doc.save()

    images = _image_blocks(_extract(pdf, tmp_path))
    raster = next((b for b in images if b.get("origin") == "embedded-raster"), None)
    assert raster is not None, (
        "no embedded-raster image block extracted from the PDF fixture "
        f"(origins seen: {[b.get('origin') for b in images]})"
    )
    assert raster.get("width") == 600
    assert raster.get("height") == 400
    frac = raster.get("page_fraction")
    assert isinstance(frac, (int, float)), "PDF embedded raster missing page_fraction"
    assert 0.0 < frac <= 1.0
