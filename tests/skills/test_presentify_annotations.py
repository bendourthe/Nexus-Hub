"""v3.15.4 Phase 3: annotated-figure overlay recreation.

Two testable surfaces:

1. The extractor captures author-added overlay shapes drawn OVER a PPTX picture
   as `annotations` metadata on that image block (image-relative bbox, text,
   fill color), and leaves a shape BESIDE the picture as normal content. Needs
   python-pptx + Pillow, so it `importorskip`s them (runs in the
   presentify-extractor workflow, skips in the deps-light ci.yml tests job).

2. The baseline builder recreates the overlay from those `annotations`: a
   registered overlay layer with one positioned region per annotation, the
   source labels, a legend, and a CSS-only view-original toggle - and it does
   NOT fabricate regions when a figure carries no annotations (the low-
   confidence / flattened degrade). The builder tests are dependency-free.

Both scripts are loaded by path via importlib, matching test_media_key_setup.py.
"""

from __future__ import annotations

import importlib.util
import json
from pathlib import Path

import pytest

_ROOT = Path(__file__).resolve().parents[2]
_BUNDLE = (
    _ROOT / "catalog" / "skills" / "specialized-domains" / "document-to-interactive-html"
)
_EXTRACT_PATH = _BUNDLE / "scripts" / "extract_content.py"
_BUILD_PATH = _BUNDLE / "scripts" / "build_presentation.py"

# A 1x1 transparent PNG, so the builder tests need no image library.
_TINY_PNG = (
    "data:image/png;base64,"
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mNk"
    "+M8AAAMBAQDJ/pLvAAAAAElFTkSuQmCC"
)


def _load(path: Path, name: str):
    spec = importlib.util.spec_from_file_location(name, path)
    assert spec and spec.loader, f"cannot load {path}"
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)
    return module


build = _load(_BUILD_PATH, "build_presentation_ann")


def _image_blocks(model: dict) -> list[dict]:
    return [
        block
        for section in model.get("sections", [])
        for block in section.get("blocks", [])
        if block.get("type") == "image"
    ]


def _annotated_deck(tmp_path: Path) -> Path:
    """A slide with a base map picture, two labeled colored region shapes over
    it, and a text box BESIDE it (which must stay normal content)."""
    from PIL import Image
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    photo = tmp_path / "map.png"
    Image.new("RGB", (800, 600), (230, 230, 230)).save(str(photo), "PNG")
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    slide.shapes.add_picture(str(photo), Inches(1), Inches(1), Inches(6), Inches(4.5))
    north = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(2), Inches(2), Inches(1.5), Inches(1)
    )
    north.fill.solid()
    north.fill.fore_color.rgb = RGBColor(0xE0, 0x40, 0x40)
    north.text_frame.text = "North zone"
    hub = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(4), Inches(2.5), Inches(1), Inches(1)
    )
    hub.fill.solid()
    hub.fill.fore_color.rgb = RGBColor(0x40, 0x80, 0xE0)
    hub.text_frame.text = "Hub"
    sidebar = slide.shapes.add_textbox(Inches(8), Inches(1), Inches(1.5), Inches(0.5))
    sidebar.text_frame.text = "Sidebar note"
    deck = tmp_path / "annotated.pptx"
    prs.save(str(deck))
    return deck


def _extract(extract, path: Path, tmp_path: Path) -> dict:
    out = tmp_path / "model.json"
    assert extract.main([str(path), "-o", str(out)]) == 0
    return json.loads(out.read_text(encoding="utf-8"))


# --- 1. extractor captures overlay shapes as annotations --------------------


def test_pptx_overlay_shapes_become_annotations(tmp_path):
    pytest.importorskip("PIL", reason="Pillow needed to build the fixture")
    pytest.importorskip("pptx", reason="python-pptx needed to build the fixture")
    extract = _load(_EXTRACT_PATH, "extract_content_ann")

    model = _extract(extract, _annotated_deck(tmp_path), tmp_path)
    images = _image_blocks(model)
    assert len(images) == 1, "expected exactly one picture block"
    annotations = images[0].get("annotations")
    assert annotations and len(annotations) == 2, (
        "the two over-map shapes should attach as annotations"
    )
    labels = {a.get("text") for a in annotations}
    assert labels == {"North zone", "Hub"}
    for ann in annotations:
        bbox = ann.get("bbox")
        assert isinstance(bbox, list) and len(bbox) == 4
        assert all(0.0 <= v <= 1.0 for v in bbox), "bbox is image-relative 0..1"
        assert ann.get("fill", "").startswith("#"), "solid fill captured as hex"
    # The shape BESIDE the picture stayed normal content, not an annotation.
    paragraphs = [
        b.get("text")
        for s in model.get("sections", [])
        for b in s.get("blocks", [])
        if b.get("type") == "paragraph"
    ]
    assert "Sidebar note" in paragraphs


# --- 2. end-to-end: annotated PPTX -> builder -> registered overlay ----------


def test_end_to_end_annotated_map_overlay(tmp_path, render_gate):
    pytest.importorskip("PIL", reason="Pillow needed to build the fixture")
    pytest.importorskip("pptx", reason="python-pptx needed to build the fixture")
    extract = _load(_EXTRACT_PATH, "extract_content_ann2")

    model_path = tmp_path / "model.json"
    assert extract.main([str(_annotated_deck(tmp_path)), "-o", str(model_path)]) == 0
    out = tmp_path / "site.html"
    assert build.main([str(model_path), "-o", str(out)]) == 0
    html = out.read_text(encoding="utf-8")
    assert 'class="fig-annotated"' in html
    assert html.count('class="fig-region"') == 2  # one per source region
    assert "North zone" in html and "Hub" in html  # source labels recreated
    assert 'class="fig-legend"' in html  # interactive legend
    assert 'class="fig-view-original"' in html  # view-original toggle


# --- 3. builder recreates the overlay from annotations (dependency-free) -----


def _annotated_model(annotations: list) -> dict:
    return {
        "schema_version": 2,
        "title": "Coverage",
        "sections": [
            {
                "heading": "Specialist coverage",
                "kind": "image",
                "blocks": [
                    {
                        "type": "image",
                        "data_uri": _TINY_PNG,
                        "alt": "Coverage map",
                        "annotations": annotations,
                    }
                ],
            }
        ],
    }


def test_builder_renders_overlay_from_annotations(tmp_path):
    annotations = [
        {"shape_type": "RECTANGLE", "bbox": [0.1, 0.2, 0.25, 0.2], "text": "North", "fill": "#E04040"},
        {"shape_type": "OVAL", "bbox": [0.5, 0.3, 0.15, 0.2], "text": "Hub", "fill": "#4080E0"},
    ]
    model = tmp_path / "m.json"
    model.write_text(json.dumps(_annotated_model(annotations)), encoding="utf-8")
    out = tmp_path / "o.html"
    assert build.main([str(model), "-o", str(out)]) == 0
    html = out.read_text(encoding="utf-8")
    assert html.count('class="fig-region"') == 2
    # Regions are positioned by image-relative percentage coordinates.
    assert "left:10.00%;top:20.00%;width:25.00%;height:20.00%" in html
    assert "--region-color:#E04040" in html
    assert "North" in html and "Hub" in html
    assert 'class="fig-legend"' in html
    assert 'class="fig-view-original"' in html
    assert "recreated from source figure" in html
    # Offline: the builder's own gate did not trip (main returned 0) and there
    # is no off-host fetch construct.
    build.assert_no_external(html)


def test_builder_rejects_malicious_fill_color(tmp_path):
    """A non-hex `fill` (an attribute-context injection attempt) is DROPPED,
    never interpolated into the `style="..."` attribute. The model is a general
    input contract, not only the trusted extractor output."""
    annotations = [
        {
            "shape_type": "RECTANGLE",
            "bbox": [0.1, 0.2, 0.25, 0.2],
            "text": "North",
            "fill": '#fff;} </style><script>alert(1)</script>',
        },
        {
            "shape_type": "OVAL",
            "bbox": [0.5, 0.3, 0.15, 0.2],
            "text": "Hub",
            "fill": '"><img src=x onerror=alert(1)>',
        },
    ]
    model = tmp_path / "m.json"
    model.write_text(json.dumps(_annotated_model(annotations)), encoding="utf-8")
    out = tmp_path / "o.html"
    assert build.main([str(model), "-o", str(out)]) == 0
    html = out.read_text(encoding="utf-8")
    assert html.count('class="fig-region"') == 2  # regions still render
    # The payloads never reach the output; both invalid fills were dropped.
    assert "<script>alert(1)</script>" not in html
    assert "onerror=alert(1)" not in html
    assert "--region-color:" not in html  # no inline color set from a bad fill


def test_builder_no_fabricated_regions_without_annotations(tmp_path):
    """The low-confidence / flattened degrade: an image with NO annotations
    renders as the plain enhanced original, with zero fabricated regions."""
    model_obj = _annotated_model([])
    model_obj["sections"][0]["blocks"][0].pop("annotations")
    model = tmp_path / "m.json"
    model.write_text(json.dumps(model_obj), encoding="utf-8")
    out = tmp_path / "o.html"
    assert build.main([str(model), "-o", str(out)]) == 0
    html = out.read_text(encoding="utf-8")
    # The template CSS always DEFINES .fig-region / .fig-overlay; assert the
    # overlay ELEMENTS were not emitted (no fabricated regions).
    assert 'class="fig-region"' not in html
    assert 'class="fig-overlay"' not in html
    assert 'class="fig-annotated"' not in html
    assert "data:image/png;base64," in html  # the base image still renders


# --- 4. headless-optional: the view-original toggle hides the overlay --------


def test_rendered_overlay_toggle(tmp_path, render_gate):
    annotations = [
        {"shape_type": "RECTANGLE", "bbox": [0.1, 0.2, 0.25, 0.2], "text": "North", "fill": "#E04040"},
    ]
    model = tmp_path / "m.json"
    model.write_text(json.dumps(_annotated_model(annotations)), encoding="utf-8")
    out = tmp_path / "o.html"
    assert build.main([str(model), "-o", str(out)]) == 0
    try:
        from playwright.sync_api import sync_playwright  # type: ignore
    except Exception:
        render_gate("no headless browser available; overlay render check skipped")
    try:
        with sync_playwright() as p:
            browser = p.chromium.launch()
            page = browser.new_page()
            page.goto(out.as_uri())
            visible_before = page.eval_on_selector(
                ".fig-overlay", "el => el.offsetParent !== null"
            )
            page.eval_on_selector(".fig-toggle", "el => el.click()")
            visible_after = page.eval_on_selector(
                ".fig-overlay", "el => getComputedStyle(el).display"
            )
            browser.close()
    except Exception:
        render_gate("headless browser present but render failed; check skipped")
    assert visible_before is True, "overlay should be visible before toggling"
    assert visible_after == "none", "view-original toggle should hide the overlay"
