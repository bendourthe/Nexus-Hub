#!/usr/bin/env python3
"""
extract_content.py - Local multi-format document extractor for the
document-to-interactive-html skill.

This script ships as a Tier-3 bundled resource: the agent invokes it via the
shell and consumes its JSON output without reading the source into context. It
maps one OR many inputs - documents (PowerPoint .pptx, Word .docx, Excel .xlsx,
PDF .pdf), source code / config, Markdown / plain text, CSV / TSV, standalone
images, or a whole directory / repository walked recursively - into the single
normalized "content model" (schema_version 2) defined in
references/content-model.md. The builder (build_presentation.py) consumes that
model and never reads a source format directly.

LOCAL-FIRST and ZERO-NETWORK by construction: every parser is a local library,
imported LAZILY inside the function that needs it so a missing library degrades
to a clear `pip install` message (and a non-zero exit) instead of a crash, and
so a missing parser for one format never blocks another. This script imports no
socket / urllib / http / requests module and opens no connection; no document
ever leaves the machine.

Parsers (lazy-imported, install only what your inputs need):
    .pptx  ->  python-pptx        (pip install python-pptx)
    .docx  ->  python-docx        (pip install python-docx)
    .xlsx  ->  openpyxl           (pip install openpyxl)
    .pdf   ->  pdfplumber         (pip install pdfplumber); pypdf is a fallback
    PDF embedded images   -> pypdf     (pip install pypdf) - optional
    PDF region/page render-> pypdfium2 (pip install pypdfium2) - optional; used
                             to rasterize vector-figure regions and scanned
                             pages. Absent => those visuals are skipped with a
                             single warning, never an error.
    scanned-page OCR      -> rapidocr-onnxruntime (pip install
                             rapidocr-onnxruntime) preferred, or pytesseract
                             when the Tesseract binary is installed - optional;
                             absent => scanned pages still ship as full-page
                             image blocks for agent-vision reading.
    images ->  Pillow (optional)  (pip install Pillow) - downscales over-budget
                                   images and converts rendered bitmaps.

Usage:
    python extract_content.py deck.pptx -o model.json
    python extract_content.py report.docx data.xlsx -o combined.json
    python extract_content.py ./inputs_folder -o model.json --max-image-bytes 1500000
    python extract_content.py ./my-repo -o repo.json --max-files 400 --max-text-bytes 300000

Output is written to the --out JSON path; all diagnostics go to stderr so stdout
stays clean. Output ordering is deterministic: sources in input order, sections
in source order, blocks in document order. A per-source `coverage` manifest
records every visual found, kept, and skipped (with reasons) so the authoring
stage can reconcile against it.
"""

from __future__ import annotations

import argparse
import base64
import hashlib
import json
import os
import re
import statistics
import sys
from pathlib import Path
from typing import NoReturn

# Extension -> language for source-code / config files (universal ingestion).
CODE_LANGUAGES = {
    ".py": "python", ".pyw": "python", ".js": "javascript",
    ".mjs": "javascript", ".cjs": "javascript", ".jsx": "jsx",
    ".ts": "typescript", ".tsx": "tsx", ".go": "go", ".rs": "rust",
    ".java": "java", ".kt": "kotlin", ".kts": "kotlin", ".c": "c",
    ".h": "c", ".cpp": "cpp", ".cc": "cpp", ".cxx": "cpp", ".hpp": "cpp",
    ".cs": "csharp", ".rb": "ruby", ".php": "php", ".swift": "swift",
    ".scala": "scala", ".r": "r", ".m": "objectivec", ".mm": "objectivec",
    ".sh": "bash", ".bash": "bash", ".zsh": "bash", ".ps1": "powershell",
    ".psm1": "powershell", ".sql": "sql", ".pl": "perl", ".lua": "lua",
    ".dart": "dart", ".vue": "vue", ".svelte": "svelte", ".json": "json",
    ".jsonc": "json", ".yaml": "yaml", ".yml": "yaml", ".toml": "toml",
    ".ini": "ini", ".cfg": "ini", ".conf": "ini", ".properties": "ini",
    ".xml": "xml", ".gradle": "groovy", ".tf": "hcl", ".proto": "protobuf",
    ".graphql": "graphql", ".css": "css", ".scss": "scss", ".less": "less",
}
MARKDOWN_EXTENSIONS = {".md", ".markdown", ".mdown", ".mkd"}
TEXT_EXTENSIONS = {".txt", ".text", ".rst", ".log"}
CSV_EXTENSIONS = {".csv", ".tsv"}
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".svg"}

EXTENSION_FORMATS = {
    ".pptx": "pptx",
    ".docx": "docx",
    ".xlsx": "xlsx",
    ".pdf": "pdf",
}
EXTENSION_FORMATS.update(dict.fromkeys(CODE_LANGUAGES, "code"))
EXTENSION_FORMATS.update(dict.fromkeys(MARKDOWN_EXTENSIONS, "markdown"))
EXTENSION_FORMATS.update(dict.fromkeys(TEXT_EXTENSIONS, "text"))
EXTENSION_FORMATS.update(dict.fromkeys(CSV_EXTENSIONS, "csv"))
EXTENSION_FORMATS.update(dict.fromkeys(IMAGE_EXTENSIONS, "image"))

# Files recognized by basename (they carry no useful suffix).
SPECIAL_BASENAMES = {
    "dockerfile": ("code", "dockerfile"),
    "makefile": ("code", "makefile"),
    "cmakelists.txt": ("code", "cmake"),
    "rakefile": ("code", "ruby"),
    "gemfile": ("code", "ruby"),
}

# Formats whose bytes are text (subject to the binary sniff during a walk).
TEXT_FORMATS = {"code", "markdown", "text", "csv"}

# Directory / repository walk controls.
DEFAULT_MAX_TEXT_BYTES = 300_000
DEFAULT_MAX_FILES = 400
CSV_MAX_ROWS = 2000
CSV_MAX_COLS = 64
IGNORE_DIRS = {
    ".git", ".hg", ".svn", ".idea", ".vscode", ".vs", "node_modules",
    ".venv", "venv", "__pycache__", "dist", "build", "out", ".next",
    ".nuxt", "target", ".pytest_cache", ".mypy_cache", ".ruff_cache",
    ".tox", ".gradle", ".terraform", "coverage", ".cache", "vendor",
    "bower_components", ".svelte-kit", ".parcel-cache", "site-packages",
}
IGNORE_BASENAMES = {
    "package-lock.json", "yarn.lock", "pnpm-lock.yaml", "poetry.lock",
    "cargo.lock", "composer.lock", "gemfile.lock", "go.sum",
}
_NUMBER_RE = re.compile(r"^-?\d+(\.\d+)?$")

DEFAULT_MAX_IMAGE_BYTES = 2_000_000

# PDF visual-extraction tuning (points unless noted).
RENDER_SCALE = 2.0  # rasterization scale for regions and scanned pages
REGION_PAD = 8.0  # cluster-merge padding
REGION_MIN_WIDTH = 60.0
REGION_MIN_HEIGHT = 40.0
REGION_MIN_PAGE_FRACTION = 0.02  # min cluster area as a fraction of page area
REGION_MAX_PAGE_FRACTION = 0.85  # single objects larger than this are backdrop
REGION_MIN_OBJECTS = 5  # min drawing objects unless curve-bearing
REGION_MIN_CURVES = 2
REGION_MAX_PER_PAGE = 6
REGION_MAX_OBJECTS = 4000  # pathological-page guard
REGION_TEXT_DENSITY_MAX = 0.20  # char-area / region-area above this = prose
CAPTION_MAX_GAP = 44.0  # max vertical gap between figure bottom and caption
LABEL_REACH = 28.0  # grow a region crop to include text this close (ticks)
CROP_MARGIN = 6.0  # breathing room around a rasterized region crop
CAPTION_MAX_LEN = 140
CAPTION_PLAIN_MAX_LEN = 90  # non-cue caption lines must be short
HEADING_SIZE_RATIO = 1.15  # heading font must beat the page median by this
HEADING_MAX_LEN = 90
SCANNED_MAX_CHARS = 32  # near-empty text layer
SCANNED_BOILERPLATE_CHARS = 200  # tolerated boilerplate on an image page
SCANNED_IMAGE_COVERAGE = 0.5
SCANNED_BOILERPLATE_COVERAGE = 0.7
OCR_LOW_CONFIDENCE = 0.80
REPEATED_ASSET_MIN_PAGES = 3  # identical image on this many pages = asset
DECK_LANDSCAPE_RATIO = 0.8  # fraction of landscape pages implying a deck PDF
DECK_MAX_AVG_CHARS = 800

_CAPTION_CUE = re.compile(
    r"^(figure|fig\.|exhibit|table|map|chart|source:)\s*", re.IGNORECASE
)
_OCR_CACHE: dict = {}
_CHART_NS = "http://schemas.openxmlformats.org/drawingml/2006/chart"


def _missing(pip_name: str) -> NoReturn:
    """Print the standard missing-dependency message and exit non-zero.

    Raises SystemExit (a BaseException, not an Exception) so it is never
    swallowed by the per-file `except Exception` guard in build_model.
    """
    print(
        f"Error: {pip_name} not installed. Please run: pip install {pip_name}",
        file=sys.stderr,
    )
    raise SystemExit(1)


# --- small shared helpers --------------------------------------------------


def _cell_str(value: object) -> str:
    """Stringify a cell value deterministically (None -> '', 3.0 -> '3')."""
    if value is None:
        return ""
    if isinstance(value, float) and value.is_integer():
        return str(int(value))
    return str(value).strip()


def _is_number(value: object) -> bool:
    """True for a real numeric cell (bool is excluded; it subclasses int)."""
    return isinstance(value, (int, float)) and not isinstance(value, bool)


def _sniff_image_type(blob: bytes) -> str:
    """Best-effort MIME type from magic bytes (default image/png)."""
    if blob.startswith(b"\xff\xd8"):
        return "image/jpeg"
    if blob.startswith(b"\x89PNG"):
        return "image/png"
    if blob.startswith((b"GIF87a", b"GIF89a")):
        return "image/gif"
    if blob.startswith(b"BM"):
        return "image/bmp"
    if blob.startswith((b"II*\x00", b"MM\x00*")):
        return "image/tiff"
    return "image/png"


def _image_dimensions(blob: bytes) -> tuple:
    """Native (width, height) in pixels via Pillow, or (None, None) when the
    library is absent or the bytes will not decode (e.g. SVG).

    Computed from the ORIGINAL blob so the signal reflects the source's native
    resolution even when the stored image is later downscaled for the budget.
    """
    try:
        from PIL import Image
    except ImportError:
        return (None, None)
    import io

    try:
        with Image.open(io.BytesIO(blob)) as image:
            return (int(image.width), int(image.height))
    except (OSError, ValueError):
        return (None, None)


def _try_downscale(blob: bytes, max_bytes: int) -> bytes | None:
    """Best-effort downscale of an over-budget image via Pillow (optional).

    Returns JPEG bytes within budget, or None if Pillow is absent or the image
    cannot be brought under budget. Never raises on a missing library.
    """
    try:
        from PIL import Image
    except ImportError:
        return None
    import io

    try:
        image = Image.open(io.BytesIO(blob)).convert("RGB")
    except (OSError, ValueError):
        return None
    scale = 1.0
    for _ in range(6):
        width, height = image.size
        new_size = (max(1, int(width * scale)), max(1, int(height * scale)))
        buffer = io.BytesIO()
        image.resize(new_size).save(buffer, format="JPEG", quality=82)
        data = buffer.getvalue()
        if len(data) <= max_bytes:
            return data
        scale *= 0.7
    return None


def _encode_image(
    blob: bytes, content_type: str, alt: str, max_bytes: int
) -> dict | None:
    """Build a base64 image block, downscaling or skipping if over budget."""
    if len(blob) > max_bytes:
        downscaled = _try_downscale(blob, max_bytes)
        if downscaled is None or len(downscaled) > max_bytes:
            print(
                f"Warning: skipping image ({len(blob)} bytes) over the "
                f"{max_bytes}-byte budget (install Pillow to downscale).",
                file=sys.stderr,
            )
            return None
        blob, content_type = downscaled, "image/jpeg"
    encoded = base64.b64encode(blob).decode("ascii")
    return {
        "type": "image",
        "data_uri": f"data:{content_type or 'image/png'};base64,{encoded}",
        "alt": alt or "Image",
    }


def _image_block(
    blob: bytes,
    content_type: str,
    alt: str,
    max_bytes: int,
    cov: dict,
    origin: str,
    page: int | None = None,
    caption: str | None = None,
    page_fraction: float | None = None,
) -> dict | None:
    """Coverage-counted image block with the schema-v2/v3 metadata fields.

    `page_fraction` (0..1) is the share of the source page/slide AREA the image
    occupies, when the caller can compute it; `width`/`height` are the native
    pixel dimensions. Both are prominence signals for the authoring stage.
    """
    cov["images_found"] += 1
    width, height = _image_dimensions(blob)
    block = _encode_image(blob, content_type, alt, max_bytes)
    if block is None:
        cov["images_skipped"] += 1
        where = f" (page {page})" if page is not None else ""
        cov["skip_reasons"].append(f"over-budget: {alt}{where}")
        return None
    block["origin"] = origin
    if page is not None:
        block["page"] = page
    if caption:
        block["caption"] = caption
    if width is not None and height is not None:
        block["width"] = width
        block["height"] = height
    if page_fraction is not None:
        block["page_fraction"] = round(max(0.0, min(1.0, float(page_fraction))), 3)
    cov["images_kept"] += 1
    return block


def _table_from_grid(grid: list) -> dict | None:
    """Turn a list-of-rows grid into a table block (first row = header)."""
    rows = [[_cell_str(cell) for cell in (row or [])] for row in grid]
    rows = [row for row in rows if any(cell for cell in row)]
    if not rows:
        return None
    return {"type": "table", "header": rows[0], "rows": rows[1:]}


def _new_coverage(path: str) -> dict:
    """Fresh per-source coverage-manifest entry (schema v2)."""
    return {
        "path": path,
        "images_found": 0,
        "images_kept": 0,
        "images_skipped": 0,
        "skip_reasons": [],
        "native_charts": 0,
        "tables": 0,
        "vector_regions_rasterized": 0,
        "vector_regions_skipped": 0,
        "scanned_pages_detected": 0,
        "ocr_pages": 0,
        "ocr_low_confidence": 0,
        "agent_read_pages": 0,
    }


# --- shared OOXML chart parsing (PPTX via python-pptx; DOCX via chart XML) ---


def _chart_hint_from_name(name: str) -> str:
    upper = name.upper()
    if "LINE" in upper or "SCATTER" in upper:
        return "line"
    if "DOUGHNUT" in upper:
        return "doughnut"
    if "PIE" in upper:
        return "pie"
    return "bar"


def _chart_xml_points(container: object) -> list:
    """Ordered c:pt values from a c:cat / c:val subtree (idx-sorted)."""
    if container is None:
        return []
    points = []
    for pt in container.iter(f"{{{_CHART_NS}}}pt"):
        value_el = pt.find(f"{{{_CHART_NS}}}v")
        text = value_el.text if value_el is not None and value_el.text else ""
        points.append((int(pt.get("idx", "0")), text))
    points.sort(key=lambda item: item[0])
    return [text for _idx, text in points]


def _chart_xml_block(xml_bytes: bytes, cov: dict, where: str) -> dict | None:
    """Parse an OOXML chart part (chart*.xml) into a chart block (stdlib only)."""
    from xml.etree import ElementTree

    c = f"{{{_CHART_NS}}}"
    try:
        root = ElementTree.fromstring(xml_bytes)
        plot_area = root.find(f".//{c}plotArea")
        if plot_area is None:
            raise ValueError("no plot area")
        chart_el = None
        hint = "bar"
        for child in plot_area:
            tag = child.tag.rsplit("}", 1)[-1]
            if tag.endswith("Chart"):
                chart_el = child
                hint = _chart_hint_from_name(tag)
                break
        if chart_el is None:
            raise ValueError("no chart element in plot area")
        categories: list = []
        series: list = []
        for number, ser in enumerate(chart_el.findall(f"{c}ser"), start=1):
            names = _chart_xml_points(ser.find(f"{c}tx"))
            cats = _chart_xml_points(ser.find(f"{c}cat"))
            raw_values = _chart_xml_points(ser.find(f"{c}val"))
            values = [float(value) if value else 0.0 for value in raw_values]
            if not values:
                continue
            if not categories and cats:
                categories = cats
            series.append(
                {"name": names[0] if names else f"Series {number}", "values": values}
            )
        if not series:
            raise ValueError("no numeric series")
        if not categories:
            categories = [str(i + 1) for i in range(len(series[0]["values"]))]
    except (ElementTree.ParseError, ValueError, TypeError) as exc:
        cov["skip_reasons"].append(f"native-chart-unreadable ({where}): {exc}")
        return None
    cov["native_charts"] += 1
    return {
        "type": "chart",
        "chart_type_hint": hint,
        "categories": categories,
        "series": series,
        "provenance": "native-chart",
    }


# --- PowerPoint ------------------------------------------------------------


def _iter_pptx_shapes(shapes: object, depth: int = 0, group_name: str | None = None):
    """Yield (leaf shape, enclosing group name or None), recursing into groups
    (depth-capped). The group name lets overlay-annotation capture group related
    shapes (a region plus its label, a legend) under one grouping."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP and depth < 8:
            name = (getattr(shape, "name", "") or "").strip() or group_name
            yield from _iter_pptx_shapes(shape.shapes, depth + 1, name)
        else:
            yield shape, group_name


def _pptx_text_block(text_frame: object) -> dict | None:
    paragraphs = [p for p in text_frame.paragraphs if p.text.strip()]
    if not paragraphs:
        return None
    if len(paragraphs) == 1 and int(getattr(paragraphs[0], "level", 0) or 0) == 0:
        return {"type": "paragraph", "text": paragraphs[0].text.strip()}
    items = [
        {"text": p.text.strip(), "depth": int(getattr(p, "level", 0) or 0)}
        for p in paragraphs
    ]
    return {"type": "bullets", "items": items}


def _pptx_table_block(table: object) -> dict | None:
    rows = [[_cell_str(cell.text) for cell in row.cells] for row in table.rows]
    if not rows:
        return None
    header: list = []
    body = rows
    try:
        if table.first_row:
            header, body = rows[0], rows[1:]
    except (AttributeError, ValueError):
        pass
    return {"type": "table", "header": header, "rows": body}


def _pptx_image_block(
    shape: object, max_bytes: int, cov: dict, slide_no: int, slide_area: float
) -> dict | None:
    try:
        image = shape.image
    except (AttributeError, ValueError):
        return None
    alt = (getattr(shape, "name", "") or "Image").strip() or "Image"
    page_fraction = None
    try:
        if slide_area and shape.width and shape.height:
            page_fraction = (float(shape.width) * float(shape.height)) / slide_area
    except (AttributeError, TypeError, ValueError):
        page_fraction = None
    return _image_block(
        image.blob,
        image.content_type,
        alt,
        max_bytes,
        cov,
        origin="shape-picture",
        page=slide_no,
        page_fraction=page_fraction,
    )


def _shape_bbox(shape: object) -> tuple | None:
    """(left, top, width, height) in EMU, or None when the geometry is absent."""
    try:
        left, top, width, height = shape.left, shape.top, shape.width, shape.height
    except (AttributeError, TypeError, ValueError):
        return None
    if None in (left, top, width, height):
        return None
    return (float(left), float(top), float(width), float(height))


def _pptx_solid_hex(fill: object) -> str | None:
    """A solid fill's fore color as #RRGGBB, or None (theme / gradient / no fill)."""
    try:
        from pptx.enum.dml import MSO_FILL

        if fill is None or fill.type != MSO_FILL.SOLID:
            return None
        return "#" + str(fill.fore_color.rgb)
    except Exception:  # noqa: BLE001 - color read is best-effort, never fatal
        return None


def _pptx_line_hex(line: object) -> str | None:
    """A line's color as #RRGGBB, or None when unset or theme-colored."""
    try:
        if line is None:
            return None
        return "#" + str(line.color.rgb)
    except Exception:  # noqa: BLE001 - color read is best-effort, never fatal
        return None


def _pptx_annotation(
    shape: object, group_name: str | None, picture_bbox: tuple
) -> dict | None:
    """Overlay-shape metadata for a shape sitting over a picture (Phase 3).

    The bbox is normalized to the underlying image's placed rectangle
    (image-relative, 0..1) - the coordinate space the overlay-recreation pattern
    consumes directly (references/figure-reconstruction.md part 5), so no
    re-projection is needed at authoring time.
    """
    bbox = _shape_bbox(shape)
    if bbox is None:
        return None
    pic_left, pic_top, pic_width, pic_height = picture_bbox
    if pic_width <= 0 or pic_height <= 0:
        return None
    annotation: dict = {
        "shape_type": getattr(shape.shape_type, "name", None) or str(shape.shape_type),
        "bbox": [
            round((bbox[0] - pic_left) / pic_width, 4),
            round((bbox[1] - pic_top) / pic_height, 4),
            round(bbox[2] / pic_width, 4),
            round(bbox[3] / pic_height, 4),
        ],
    }
    if getattr(shape, "has_text_frame", False):
        text = shape.text_frame.text.strip()
        if text:
            annotation["text"] = text
    fill = _pptx_solid_hex(getattr(shape, "fill", None))
    if fill:
        annotation["fill"] = fill
    line = _pptx_line_hex(getattr(shape, "line", None))
    if line:
        annotation["line"] = line
    if group_name:
        annotation["group"] = group_name
    return annotation


def _pptx_chart_block(shape: object, cov: dict, slide_no: int) -> dict | None:
    """Native PPTX chart -> chart block with the source's real series values."""
    try:
        chart = shape.chart
        plot = chart.plots[0]
        categories = [str(category) for category in plot.categories]
        series: list = []
        none_seen = False
        for number, ser in enumerate(chart.series, start=1):
            try:
                name = str(ser.name)
            except (AttributeError, KeyError, ValueError):
                name = f"Series {number}"
            values: list = []
            for value in ser.values:
                if value is None:
                    none_seen = True
                    values.append(0.0)
                else:
                    values.append(float(value))
            if values:
                series.append({"name": name, "values": values})
        if not series or not categories:
            raise ValueError("empty chart data")
        try:
            hint = _chart_hint_from_name(str(chart.chart_type))
        except (AttributeError, ValueError):
            hint = "bar"
        caption = ""
        if getattr(chart, "has_title", False):
            caption = chart.chart_title.text_frame.text.strip()
    except Exception as exc:  # noqa: BLE001 - logged to coverage, never fatal
        cov["skip_reasons"].append(f"native-chart-unreadable (slide {slide_no}): {exc}")
        return None
    if none_seen:
        print(
            f"Warning: chart on slide {slide_no} has empty data points "
            "(recorded as 0.0).",
            file=sys.stderr,
        )
    cov["native_charts"] += 1
    block = {
        "type": "chart",
        "chart_type_hint": hint,
        "categories": categories,
        "series": series,
        "provenance": "native-chart",
    }
    if caption:
        block["caption"] = caption
    return block


def _pptx_kind(index: int, heading: str, blocks: list) -> str:
    if index == 0:
        return "title"
    has_content = any(block["type"] != "notes" for block in blocks)
    if heading and not has_content:
        return "section-break"
    return "content"


def _extract_pptx(path: str, max_bytes: int, cov: dict) -> tuple[str, list]:
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
    except ImportError:
        _missing("python-pptx")

    presentation = Presentation(path)
    slide_area = float(presentation.slide_width or 0) * float(
        presentation.slide_height or 0
    )
    doc_title: str | None = None
    sections: list = []
    for index, slide in enumerate(presentation.slides):
        slide_no = index + 1
        title_shape = slide.shapes.title
        heading = title_shape.text.strip() if title_shape is not None else ""
        if doc_title is None and heading:
            doc_title = heading
        skip_ids = {title_shape.shape_id} if title_shape is not None else set()
        subheading: str | None = None
        for placeholder in slide.placeholders:
            try:
                if placeholder.placeholder_format.type == PP_PLACEHOLDER.SUBTITLE:
                    subheading = placeholder.text.strip()
                    skip_ids.add(placeholder.shape_id)
                    break
            except (AttributeError, ValueError):
                continue
        shape_items = [
            (shape, group)
            for shape, group in _iter_pptx_shapes(slide.shapes)
            if shape.shape_id not in skip_ids
        ]
        # Pass 1: picture placed rectangles, for overlay-annotation assignment.
        picture_bboxes: dict = {}  # shape index -> bbox (EMU)
        for index, (shape, _group) in enumerate(shape_items):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                bbox = _shape_bbox(shape)
                if bbox is not None:
                    picture_bboxes[index] = bbox
        # Pass 2: a non-picture shape whose CENTER lies inside a picture, and
        # which is smaller than that picture, is an author-added overlay - it
        # attaches to that image as annotation metadata and is consumed (not
        # also emitted as a stray text block). Tables and charts never qualify.
        annotations_for: dict = {}  # picture index -> [annotation, ...]
        consumed: set = set()
        for index, (shape, group) in enumerate(shape_items):
            if index in picture_bboxes:
                continue
            if getattr(shape, "has_table", False) or getattr(shape, "has_chart", False):
                continue
            bbox = _shape_bbox(shape)
            if bbox is None:
                continue
            center_x, center_y = bbox[0] + bbox[2] / 2, bbox[1] + bbox[3] / 2
            for pic_index, pic_bbox in picture_bboxes.items():
                inside = (
                    pic_bbox[0] <= center_x <= pic_bbox[0] + pic_bbox[2]
                    and pic_bbox[1] <= center_y <= pic_bbox[1] + pic_bbox[3]
                )
                if inside and bbox[2] * bbox[3] < pic_bbox[2] * pic_bbox[3]:
                    annotation = _pptx_annotation(shape, group, pic_bbox)
                    if annotation is not None:
                        annotations_for.setdefault(pic_index, []).append(annotation)
                        consumed.add(index)
                    break
        # Pass 3: build blocks in document order, attaching annotations to their
        # picture's image block.
        blocks: list = []
        for index, (shape, _group) in enumerate(shape_items):
            if index in consumed:
                continue
            block = None
            if getattr(shape, "has_table", False):
                block = _pptx_table_block(shape.table)
                if block is not None:
                    cov["tables"] += 1
            elif getattr(shape, "has_chart", False):
                block = _pptx_chart_block(shape, cov, slide_no)
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                block = _pptx_image_block(shape, max_bytes, cov, slide_no, slide_area)
                if block is not None and annotations_for.get(index):
                    block["annotations"] = annotations_for[index]
            elif getattr(shape, "has_text_frame", False):
                block = _pptx_text_block(shape.text_frame)
            if block is not None:
                blocks.append(block)
        if slide.has_notes_slide:
            note = slide.notes_slide.notes_text_frame.text.strip()
            if note:
                blocks.append({"type": "notes", "text": note})
        sections.append(
            {
                "heading": heading,
                "subheading": subheading,
                "kind": _pptx_kind(index, heading, blocks),
                "source_index": 0,
                "blocks": blocks,
            }
        )
    return doc_title or Path(path).stem, sections


# --- Word ------------------------------------------------------------------


def _iter_docx_blocks(document: object):
    """Yield Paragraph and Table objects in document order."""
    from docx.oxml.table import CT_Tbl
    from docx.oxml.text.paragraph import CT_P
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    for child in document.element.body.iterchildren():
        if isinstance(child, CT_P):
            yield Paragraph(child, document)
        elif isinstance(child, CT_Tbl):
            yield Table(child, document)


def _docx_list_level(paragraph: object) -> int | None:
    """Return a list nesting level (0-based) or None when not a list item."""
    from docx.oxml.ns import qn

    properties = paragraph._p.pPr
    if properties is None:
        name = (paragraph.style.name if paragraph.style else "") or ""
        return 0 if name.startswith("List") else None
    num_pr = properties.find(qn("w:numPr"))
    if num_pr is None:
        name = (paragraph.style.name if paragraph.style else "") or ""
        return 0 if name.startswith("List") else None
    level = num_pr.find(qn("w:ilvl"))
    if level is not None:
        return int(level.get(qn("w:val")) or 0)
    return 0


def _docx_paragraph_images(
    paragraph: object, document: object, max_bytes: int, cov: dict
) -> list:
    from docx.oxml.ns import qn

    blocks: list = []
    for blip in paragraph._p.findall(".//" + qn("a:blip")):
        rid = blip.get(qn("r:embed"))
        if not rid:
            continue
        try:
            part = document.part.related_parts[rid]
        except KeyError:
            continue
        block = _image_block(
            part.blob,
            getattr(part, "content_type", "image/png"),
            "Image",
            max_bytes,
            cov,
            origin="inline-image",
        )
        if block is not None:
            blocks.append(block)
    return blocks


def _docx_paragraph_charts(paragraph: object, document: object, cov: dict) -> list:
    """Native DOCX chart parts referenced from this paragraph's drawings."""
    from docx.oxml.ns import qn

    blocks: list = []
    for chart_ref in paragraph._p.findall(f".//{{{_CHART_NS}}}chart"):
        rid = chart_ref.get(qn("r:id"))
        if not rid:
            continue
        try:
            part = document.part.related_parts[rid]
        except KeyError:
            cov["skip_reasons"].append(f"native-chart-unreadable (docx rel {rid})")
            continue
        block = _chart_xml_block(part.blob, cov, where=f"docx rel {rid}")
        if block is not None:
            blocks.append(block)
    return blocks


def _docx_table_block(table: object) -> dict | None:
    rows = [[_cell_str(cell.text) for cell in row.cells] for row in table.rows]
    return _table_from_grid(rows)


def _extract_docx(path: str, max_bytes: int, cov: dict) -> tuple[str, list]:
    try:
        from docx import Document
    except ImportError:
        _missing("python-docx")

    document = Document(path)
    doc_title: str | None = None
    sections: list = []
    current: dict | None = None
    pending: list = []

    def flush_bullets() -> None:
        if pending and current is not None:
            current["blocks"].append({"type": "bullets", "items": list(pending)})
        pending.clear()

    def ensure_section() -> dict:
        nonlocal current
        if current is None:
            current = {
                "heading": "Overview",
                "subheading": None,
                "kind": "content",
                "source_index": 0,
                "blocks": [],
            }
            sections.append(current)
        return current

    for item in _iter_docx_blocks(document):
        if hasattr(item, "rows"):  # Table
            flush_bullets()
            block = _docx_table_block(item)
            if block is not None:
                cov["tables"] += 1
                ensure_section()["blocks"].append(block)
            continue
        style = (item.style.name if item.style else "") or ""
        text = item.text.strip()
        if style == "Title":
            if not doc_title:
                doc_title = text
            continue
        if style.startswith("Heading"):
            flush_bullets()
            current = {
                "heading": text,
                "subheading": None,
                "kind": "content",
                "source_index": 0,
                "blocks": [],
            }
            sections.append(current)
            continue
        visuals = _docx_paragraph_images(item, document, max_bytes, cov)
        visuals.extend(_docx_paragraph_charts(item, document, cov))
        if visuals:
            flush_bullets()
            ensure_section()["blocks"].extend(visuals)
        if not text:
            continue
        level = _docx_list_level(item)
        if level is not None:
            ensure_section()
            pending.append({"text": text, "depth": level})
        else:
            flush_bullets()
            ensure_section()["blocks"].append({"type": "paragraph", "text": text})
    flush_bullets()
    return doc_title or Path(path).stem, sections


# --- Excel -----------------------------------------------------------------


def _trim_grid(rows: list) -> list:
    rows = [
        list(row)
        for row in rows
        if any(cell is not None and str(cell).strip() != "" for cell in row)
    ]
    if not rows:
        return []
    width = max(len(row) for row in rows)
    rows = [row + [None] * (width - len(row)) for row in rows]
    keep = [
        col
        for col in range(width)
        if any(row[col] is not None and str(row[col]).strip() != "" for row in rows)
    ]
    if not keep:
        return []
    return [[row[col] for col in keep] for row in rows]


def _chart_hint(category_count: int, series_count: int) -> str:
    if series_count == 1 and category_count <= 6:
        return "pie"
    if category_count > 12:
        return "line"
    return "bar"


def _grid_to_block(grid: list) -> dict | None:
    if len(grid) >= 2 and len(grid[0]) >= 2:
        header = [_cell_str(cell) for cell in grid[0]]
        body = grid[1:]
        categories = [_cell_str(row[0]) for row in body]
        series: list = []
        for col in range(1, len(header)):
            column = [row[col] for row in body]
            if column and all(_is_number(value) for value in column):
                series.append(
                    {
                        "name": header[col] or f"Series {col}",
                        "values": [float(value) for value in column],
                    }
                )
        if series and any(categories):
            return {
                "type": "chart",
                "chart_type_hint": _chart_hint(len(categories), len(series)),
                "categories": categories,
                "series": series,
                "provenance": "source-data",
            }
    return _table_from_grid(grid)


def _extract_xlsx(path: str, cov: dict) -> tuple[str, list]:
    try:
        from openpyxl import load_workbook
    except ImportError:
        _missing("openpyxl")

    workbook = load_workbook(path, data_only=True, read_only=True)
    sections: list = []
    for worksheet in workbook.worksheets:
        grid = _trim_grid([list(row) for row in worksheet.iter_rows(values_only=True)])
        if not grid:
            continue
        block = _grid_to_block(grid)
        if block is not None and block["type"] == "table":
            cov["tables"] += 1
        sections.append(
            {
                "heading": worksheet.title,
                "subheading": None,
                "kind": "data",
                "source_index": 0,
                "blocks": [block] if block is not None else [],
            }
        )
    workbook.close()
    return Path(path).stem, sections


# --- PDF -------------------------------------------------------------------


def _pdf_page_section(
    text: str, page_number: int, heading_override: str | None = None
) -> dict:
    """Split page text into paragraph blocks under a detected heading."""
    lines = text.splitlines()
    heading = f"Page {page_number}"
    start = 0
    if heading_override:
        heading = heading_override
        # Drop the first occurrence of the heading line from the body text.
        stripped = heading_override.strip()
        lines = list(lines)
        for index, line in enumerate(lines):
            if line.strip() == stripped:
                del lines[index]
                break
    else:
        for index, line in enumerate(lines):
            if line.strip():
                if len(line.strip()) <= 80:
                    heading, start = line.strip(), index + 1
                else:
                    start = index
                break
    paragraphs: list = []
    current: list = []
    for line in lines[start:]:
        if line.strip():
            current.append(line.strip())
        elif current:
            paragraphs.append(" ".join(current))
            current = []
    if current:
        paragraphs.append(" ".join(current))
    blocks = [{"type": "paragraph", "text": para} for para in paragraphs if para]
    return {
        "heading": heading,
        "subheading": None,
        "kind": "content",
        "source_index": 0,
        "blocks": blocks,
    }


def _pdf_collect_rasters(path: str, cov: dict) -> dict | None:
    """pypdf pass over the PDF: {page_index: [{name, blob, hash}]}.

    Returns None when pypdf is absent (embedded-raster extraction disabled).
    """
    try:
        from pypdf import PdfReader
    except ImportError:
        print(
            "Warning: pypdf not installed; skipping embedded PDF image "
            "extraction. Please run: pip install pypdf",
            file=sys.stderr,
        )
        cov["skip_reasons"].append("pdf-images-unavailable: pypdf not installed")
        return None

    result: dict = {}
    try:
        reader = PdfReader(path)
    except Exception as exc:  # noqa: BLE001 - logged, images skipped
        cov["skip_reasons"].append(f"pdf-images-unavailable: {exc}")
        return None
    for index, page in enumerate(reader.pages):
        entries: list = []
        try:
            images = list(page.images)
        except Exception as exc:  # noqa: BLE001 - per-page guard
            cov["skip_reasons"].append(f"page {index + 1} image read failed: {exc}")
            images = []
        for image in images:
            try:
                blob = image.data
            except Exception as exc:  # noqa: BLE001 - per-image guard
                cov["images_found"] += 1
                cov["images_skipped"] += 1
                cov["skip_reasons"].append(
                    f"undecodable-image (page {index + 1}): {exc}"
                )
                continue
            entries.append(
                {
                    "name": getattr(image, "name", "") or "Image",
                    "blob": blob,
                    "hash": hashlib.sha256(blob).hexdigest(),
                }
            )
        result[index] = entries
    return result


def _repeated_hashes(raster_map: dict) -> dict:
    """Hashes of images repeated across >= REPEATED_ASSET_MIN_PAGES pages."""
    pages_by_hash: dict = {}
    for index, entries in raster_map.items():
        for entry in entries:
            pages_by_hash.setdefault(entry["hash"], set()).add(index)
    return {
        digest: pages
        for digest, pages in pages_by_hash.items()
        if len(pages) >= REPEATED_ASSET_MIN_PAGES
    }


def _pdf_open_renderer(path: str) -> object | None:
    """Open the optional pypdfium2 renderer (None when unavailable)."""
    try:
        import pypdfium2 as pdfium
    except ImportError:
        return None
    try:
        from PIL import Image  # noqa: F401 - required by bitmap.to_pil()
    except ImportError:
        return None
    try:
        return pdfium.PdfDocument(path)
    except Exception as exc:  # noqa: BLE001 - renderer is optional
        print(f"Warning: could not open PDF for rendering: {exc}", file=sys.stderr)
        return None


def _pdf_render_region(
    renderer: object, page_index: int, bbox: tuple | None
) -> bytes | None:
    """Rasterize a page (or a page region) to PNG bytes via pypdfium2."""
    import io

    try:
        page = renderer[page_index]
        bitmap = page.render(scale=RENDER_SCALE)
        pil = bitmap.to_pil()
        if bbox is not None:
            x0, top, x1, bottom = bbox
            pil = pil.crop(
                (
                    max(0, int(x0 * RENDER_SCALE)),
                    max(0, int(top * RENDER_SCALE)),
                    min(pil.width, int(x1 * RENDER_SCALE)),
                    min(pil.height, int(bottom * RENDER_SCALE)),
                )
            )
        buffer = io.BytesIO()
        pil.save(buffer, format="PNG")
        return buffer.getvalue()
    except Exception as exc:  # noqa: BLE001 - rendering is best-effort
        print(f"Warning: page {page_index + 1} render failed: {exc}", file=sys.stderr)
        return None


def _bbox_overlap(a: tuple, b: tuple) -> float:
    """Intersection area of two (x0, top, x1, bottom) boxes."""
    width = min(a[2], b[2]) - max(a[0], b[0])
    height = min(a[3], b[3]) - max(a[1], b[1])
    if width <= 0 or height <= 0:
        return 0.0
    return width * height


def _pdf_figure_regions(page: object, table_bboxes: list, image_bboxes: list) -> list:
    """Cluster vector drawing objects into candidate figure-region bboxes."""
    page_area = float(page.width) * float(page.height) or 1.0
    objects: list = []
    for kind in ("curves", "rects", "lines"):
        for obj in getattr(page, kind, None) or []:
            bbox = (obj["x0"], obj["top"], obj["x1"], obj["bottom"])
            width = bbox[2] - bbox[0]
            height = bbox[3] - bbox[1]
            if width * height / page_area > REGION_MAX_PAGE_FRACTION:
                continue  # full-page backdrop / border
            objects.append((bbox, kind))
    if not objects or len(objects) > REGION_MAX_OBJECTS:
        return []

    clusters: list = []  # each: [x0, top, x1, bottom, count, curve_count]
    for bbox, kind in objects:
        padded = (
            bbox[0] - REGION_PAD,
            bbox[1] - REGION_PAD,
            bbox[2] + REGION_PAD,
            bbox[3] + REGION_PAD,
        )
        merged = None
        for cluster in clusters:
            if _bbox_overlap(padded, tuple(cluster[:4])) > 0:
                cluster[0] = min(cluster[0], bbox[0])
                cluster[1] = min(cluster[1], bbox[1])
                cluster[2] = max(cluster[2], bbox[2])
                cluster[3] = max(cluster[3], bbox[3])
                cluster[4] += 1
                cluster[5] += 1 if kind == "curves" else 0
                merged = cluster
                break
        if merged is None:
            clusters.append(
                [bbox[0], bbox[1], bbox[2], bbox[3], 1, 1 if kind == "curves" else 0]
            )
    # Merge overlapping clusters until stable (bounded passes).
    for _ in range(8):
        changed = False
        result: list = []
        for cluster in clusters:
            merged = None
            for kept in result:
                pad_kept = (
                    kept[0] - REGION_PAD,
                    kept[1] - REGION_PAD,
                    kept[2] + REGION_PAD,
                    kept[3] + REGION_PAD,
                )
                if _bbox_overlap(pad_kept, tuple(cluster[:4])) > 0:
                    kept[0] = min(kept[0], cluster[0])
                    kept[1] = min(kept[1], cluster[1])
                    kept[2] = max(kept[2], cluster[2])
                    kept[3] = max(kept[3], cluster[3])
                    kept[4] += cluster[4]
                    kept[5] += cluster[5]
                    merged = kept
                    changed = True
                    break
            if merged is None:
                result.append(cluster)
        clusters = result
        if not changed:
            break

    chars = getattr(page, "chars", None) or []
    regions: list = []
    for cluster in clusters:
        bbox = tuple(cluster[:4])
        width = bbox[2] - bbox[0]
        height = bbox[3] - bbox[1]
        if width < REGION_MIN_WIDTH or height < REGION_MIN_HEIGHT:
            continue
        if width * height / page_area < REGION_MIN_PAGE_FRACTION:
            continue
        if cluster[4] < REGION_MIN_OBJECTS and cluster[5] < REGION_MIN_CURVES:
            continue
        region_area = width * height
        if any(
            _bbox_overlap(bbox, table) / region_area >= 0.5 for table in table_bboxes
        ):
            continue  # detected table, already extracted as a table block
        if any(
            _bbox_overlap(bbox, image) / region_area >= 0.8 for image in image_bboxes
        ):
            continue  # covered by an embedded raster we already extracted
        char_area = sum(
            (char["x1"] - char["x0"]) * (char["bottom"] - char["top"])
            for char in chars
            if bbox[0] <= (char["x0"] + char["x1"]) / 2 <= bbox[2]
            and bbox[1] <= (char["top"] + char["bottom"]) / 2 <= bbox[3]
        )
        if char_area / region_area > REGION_TEXT_DENSITY_MAX:
            continue  # dense prose, not a figure
        regions.append(bbox)
    regions.sort(key=lambda item: (item[1], item[0]))
    return regions[:REGION_MAX_PER_PAGE]


def _expand_region_with_labels(page: object, bbox: tuple) -> tuple:
    """Grow a figure-region crop to include nearby label text (axis ticks,
    category labels, captions) so the rasterized figure stays readable.

    Single-pass by design: only text near the ORIGINAL bbox is pulled in, so
    the crop cannot creep across the page.
    """
    x0, top, x1, bottom = bbox
    for char in getattr(page, "chars", None) or []:
        center_x = (char["x0"] + char["x1"]) / 2
        center_y = (char["top"] + char["bottom"]) / 2
        if (
            bbox[0] - LABEL_REACH <= center_x <= bbox[2] + LABEL_REACH
            and bbox[1] - LABEL_REACH <= center_y <= bbox[3] + LABEL_REACH
        ):
            x0 = min(x0, char["x0"])
            top = min(top, char["top"])
            x1 = max(x1, char["x1"])
            bottom = max(bottom, char["bottom"])
    return (
        max(0.0, x0 - CROP_MARGIN),
        max(0.0, top - CROP_MARGIN),
        min(float(page.width), x1 + CROP_MARGIN),
        min(float(page.height), bottom + CROP_MARGIN),
    )


def _pdf_text_lines(page: object) -> list:
    """Positioned text lines (empty when the pdfplumber API is too old)."""
    try:
        lines = page.extract_text_lines()
    except (AttributeError, TypeError):
        return []
    result: list = []
    for line in lines or []:
        text = (line.get("text") or "").strip()
        if text:
            result.append(
                {
                    "text": text,
                    "x0": line["x0"],
                    "x1": line["x1"],
                    "top": line["top"],
                    "bottom": line["bottom"],
                }
            )
    return result


def _match_caption(bbox: tuple, lines: list) -> str | None:
    """The caption line sitting directly below a figure bbox, if any."""
    x0, _top, x1, bottom = bbox
    best: str | None = None
    best_key: tuple | None = None
    for line in lines:
        if line["top"] < bottom - 2 or line["top"] > bottom + CAPTION_MAX_GAP:
            continue
        if line["top"] - bottom < 6 and line["x0"] >= x0 and line["x1"] <= x1:
            continue  # axis-hugging tick-label row, not a caption
        overlap = min(x1, line["x1"]) - max(x0, line["x0"])
        if overlap <= 0:
            continue
        ratio = overlap / max(1.0, min(x1 - x0, line["x1"] - line["x0"]))
        if ratio < 0.3:
            continue
        text = line["text"]
        if len(text) > CAPTION_MAX_LEN:
            continue
        cue = bool(_CAPTION_CUE.match(text))
        if not cue and len(text) > CAPTION_PLAIN_MAX_LEN:
            continue
        key = (0 if cue else 1, line["top"] - bottom)
        if best_key is None or key < best_key:
            best, best_key = text, key
    return best


def _pdf_heading(page: object) -> str | None:
    """Typographic heading: the page's largest-font short line near the top."""
    try:
        words = page.extract_words(extra_attrs=["size"])
    except (AttributeError, TypeError):
        return None
    if not words:
        return None
    sizes = [word["size"] for word in words if word.get("size")]
    if not sizes:
        return None
    median_size = statistics.median(sizes)
    lines: dict = {}
    for word in words:
        key = round(word["top"] / 3.0)
        lines.setdefault(key, []).append(word)
    candidates: list = []
    for key in sorted(lines):
        group = sorted(lines[key], key=lambda word: word["x0"])
        text = " ".join(word["text"] for word in group).strip()
        if not text or len(text) > HEADING_MAX_LEN:
            continue
        top = min(word["top"] for word in group)
        if top > float(page.height) * 0.45:
            continue
        avg_size = sum(word.get("size") or 0 for word in group) / len(group)
        if avg_size >= median_size * HEADING_SIZE_RATIO:
            candidates.append((avg_size, -top, text))
    if not candidates:
        return None
    candidates.sort(reverse=True)
    return candidates[0][2]


def _pdf_page_is_scanned(text: str, page: object) -> bool:
    """True when the page is image-dominated with a near-empty text layer."""
    images = getattr(page, "images", None) or []
    if not images:
        return False
    page_area = float(page.width) * float(page.height) or 1.0
    coverage = (
        sum(
            max(0.0, image["x1"] - image["x0"])
            * max(0.0, image["bottom"] - image["top"])
            for image in images
        )
        / page_area
    )
    chars = len(text.strip())
    if chars < SCANNED_MAX_CHARS and coverage >= SCANNED_IMAGE_COVERAGE:
        return True
    return chars < SCANNED_BOILERPLATE_CHARS and coverage >= (
        SCANNED_BOILERPLATE_COVERAGE
    )


def _load_ocr_engine() -> tuple | None:
    """Load the optional local OCR engine once per run (None when absent)."""
    if "engine" in _OCR_CACHE:
        return _OCR_CACHE["engine"]
    engine: tuple | None = None
    try:
        from rapidocr_onnxruntime import RapidOCR

        engine = ("rapidocr", RapidOCR())
    except ImportError:
        try:
            import pytesseract

            try:
                pytesseract.get_tesseract_version()
                engine = ("tesseract", pytesseract)
            except Exception:  # noqa: BLE001 - binary missing => unavailable
                engine = None
        except ImportError:
            engine = None
    _OCR_CACHE["engine"] = engine
    return engine


def _ocr_entries_to_blocks(entries: list, cov: dict) -> list:
    """Group positioned OCR entries into paragraph/table blocks.

    Entries: {x0, x1, top, bottom, text, score}. Rows are built from vertical
    overlap; runs of rows with the same multi-cell shape become a table block,
    everything else merges into paragraphs by vertical proximity.
    """
    if not entries:
        return []
    entries = sorted(entries, key=lambda e: (e["top"], e["x0"]))
    rows: list = []
    for entry in entries:
        placed = False
        for row in rows:
            overlap = min(row["bottom"], entry["bottom"]) - max(
                row["top"], entry["top"]
            )
            height = min(row["bottom"] - row["top"], entry["bottom"] - entry["top"])
            if height > 0 and overlap / height >= 0.5:
                row["cells"].append(entry)
                row["top"] = min(row["top"], entry["top"])
                row["bottom"] = max(row["bottom"], entry["bottom"])
                placed = True
                break
        if not placed:
            rows.append(
                {"top": entry["top"], "bottom": entry["bottom"], "cells": [entry]}
            )
    for row in rows:
        row["cells"].sort(key=lambda e: e["x0"])
    rows.sort(key=lambda row: row["top"])

    heights = [row["bottom"] - row["top"] for row in rows]
    median_height = statistics.median(heights) if heights else 12.0

    def finish_paragraph(chunk: list, blocks: list) -> None:
        if not chunk:
            return
        text = " ".join(
            " ".join(cell["text"] for cell in row["cells"]) for row in chunk
        ).strip()
        if not text:
            return
        score = min(cell["score"] for row in chunk for cell in row["cells"])
        if score < OCR_LOW_CONFIDENCE:
            cov["ocr_low_confidence"] += 1
        blocks.append(
            {
                "type": "paragraph",
                "text": text,
                "provenance": "ocr",
                "ocr_confidence": round(score, 3),
            }
        )

    def finish_table(chunk: list, blocks: list) -> None:
        grid = [[cell["text"] for cell in row["cells"]] for row in chunk]
        block = _table_from_grid(grid)
        if block is None:
            return
        score = min(cell["score"] for row in chunk for cell in row["cells"])
        if score < OCR_LOW_CONFIDENCE:
            cov["ocr_low_confidence"] += 1
        block["provenance"] = "ocr"
        block["ocr_confidence"] = round(score, 3)
        cov["tables"] += 1
        blocks.append(block)

    blocks: list = []
    paragraph_chunk: list = []
    table_chunk: list = []
    previous_bottom: float | None = None
    for row in rows:
        cell_count = len(row["cells"])
        if cell_count >= 2:
            if table_chunk and len(table_chunk[-1]["cells"]) != cell_count:
                if len(table_chunk) >= 2:
                    finish_paragraph(paragraph_chunk, blocks)
                    paragraph_chunk = []
                    finish_table(table_chunk, blocks)
                else:
                    paragraph_chunk.extend(table_chunk)
                table_chunk = []
            table_chunk.append(row)
        else:
            if len(table_chunk) >= 2:
                finish_paragraph(paragraph_chunk, blocks)
                paragraph_chunk = []
                finish_table(table_chunk, blocks)
            else:
                paragraph_chunk.extend(table_chunk)
            table_chunk = []
            gap_break = (
                previous_bottom is not None
                and row["top"] - previous_bottom > 1.8 * median_height
            )
            if gap_break:
                finish_paragraph(paragraph_chunk, blocks)
                paragraph_chunk = []
            paragraph_chunk.append(row)
        previous_bottom = row["bottom"]
    if len(table_chunk) >= 2:
        finish_paragraph(paragraph_chunk, blocks)
        paragraph_chunk = []
        finish_table(table_chunk, blocks)
    else:
        paragraph_chunk.extend(table_chunk)
    finish_paragraph(paragraph_chunk, blocks)
    return blocks


def _ocr_rapidocr(engine: object, png: bytes, cov: dict) -> list:
    import io

    import numpy
    from PIL import Image

    array = numpy.asarray(Image.open(io.BytesIO(png)).convert("RGB"))
    try:
        result, _elapsed = engine(array)
    except Exception as exc:  # noqa: BLE001 - OCR is best-effort
        print(f"Warning: OCR failed: {exc}", file=sys.stderr)
        return []
    entries: list = []
    for box, text, score in result or []:
        text = (text or "").strip()
        if not text:
            continue
        entries.append(
            {
                "x0": min(point[0] for point in box),
                "x1": max(point[0] for point in box),
                "top": min(point[1] for point in box),
                "bottom": max(point[1] for point in box),
                "text": text,
                "score": float(score),
            }
        )
    return _ocr_entries_to_blocks(entries, cov)


def _ocr_tesseract(module: object, png: bytes, cov: dict) -> list:
    import io

    from PIL import Image

    try:
        data = module.image_to_data(
            Image.open(io.BytesIO(png)), output_type=module.Output.DICT
        )
    except Exception as exc:  # noqa: BLE001 - OCR is best-effort
        print(f"Warning: OCR failed: {exc}", file=sys.stderr)
        return []
    entries: list = []
    count = len(data.get("text", []))
    for index in range(count):
        text = (data["text"][index] or "").strip()
        conf = float(data["conf"][index])
        if not text or conf < 0:
            continue
        left = float(data["left"][index])
        top = float(data["top"][index])
        entries.append(
            {
                "x0": left,
                "x1": left + float(data["width"][index]),
                "top": top,
                "bottom": top + float(data["height"][index]),
                "text": text,
                "score": conf / 100.0,
            }
        )
    return _ocr_entries_to_blocks(entries, cov)


def _ocr_page(png: bytes, cov: dict) -> list:
    engine = _load_ocr_engine()
    if engine is None:
        return []
    kind, impl = engine
    if kind == "rapidocr":
        return _ocr_rapidocr(impl, png, cov)
    return _ocr_tesseract(impl, png, cov)


def _pdf_scanned_section(
    page_number: int,
    page_index: int,
    renderer: object | None,
    rasters: list,
    max_bytes: int,
    cov: dict,
) -> dict | None:
    """Scanned page -> OCR text blocks (tier A) + a full-page image (tier B)."""
    cov["scanned_pages_detected"] += 1
    page_png: bytes | None = None
    content_type = "image/png"
    if renderer is not None:
        page_png = _pdf_render_region(renderer, page_index, None)
    if page_png is None and rasters:
        largest = max(rasters, key=lambda entry: len(entry["blob"]))
        page_png = largest["blob"]
        content_type = _sniff_image_type(page_png)
    if page_png is None:
        cov["skip_reasons"].append(
            f"scanned-page-unreadable (page {page_number}): install pypdfium2"
        )
        return None
    blocks = _ocr_page(page_png, cov)
    if blocks:
        cov["ocr_pages"] += 1
    else:
        cov["agent_read_pages"] += 1
    image = _image_block(
        page_png,
        content_type,
        f"Scanned page {page_number}",
        max_bytes,
        cov,
        origin="scanned-page",
        page=page_number,
    )
    if image is not None:
        blocks.append(image)
    return {
        "heading": f"Page {page_number}",
        "subheading": None,
        "kind": "content",
        "source_index": 0,
        "blocks": blocks,
    }


def _extract_pdf_pypdf(path: str, max_bytes: int, cov: dict) -> tuple[str, list]:
    """Text + embedded-raster fallback when pdfplumber is not installed.

    No layout metadata is available here, so vector-figure regions, captions,
    typographic headings, and scanned-page detection are pdfplumber-only.
    """
    try:
        from pypdf import PdfReader
    except ImportError:
        _missing("pdfplumber")

    raster_map = _pdf_collect_rasters(path, cov) or {}
    repeated = _repeated_hashes(raster_map)
    emitted_hashes: set = set()
    reader = PdfReader(path)
    doc_title: str | None = None
    sections: list = []
    for index, page in enumerate(reader.pages):
        section = _pdf_page_section(page.extract_text() or "", index + 1)
        for entry in raster_map.get(index, []):
            if entry["hash"] in repeated:
                if entry["hash"] in emitted_hashes:
                    cov["images_found"] += 1
                    cov["images_skipped"] += 1
                    continue
                emitted_hashes.add(entry["hash"])
                pages = sorted(repeated[entry["hash"]])
                cov["skip_reasons"].append(
                    f"repeated-asset: image on {len(pages)} pages kept once "
                    f"(page {pages[0] + 1})"
                )
            block = _image_block(
                entry["blob"],
                _sniff_image_type(entry["blob"]),
                entry["name"],
                max_bytes,
                cov,
                origin="embedded-raster",
                page=index + 1,
            )
            if block is not None:
                section["blocks"].append(block)
        if doc_title is None and not section["heading"].startswith("Page "):
            doc_title = section["heading"]
        sections.append(section)
    return doc_title or Path(path).stem, sections


def _extract_pdf(path: str, max_bytes: int, cov: dict) -> tuple[str, list]:
    try:
        import pdfplumber
    except ImportError:
        return _extract_pdf_pypdf(path, max_bytes, cov)

    raster_map = _pdf_collect_rasters(path, cov)
    repeated = _repeated_hashes(raster_map) if raster_map else {}
    emitted_repeats: set = set()
    renderer = _pdf_open_renderer(path)
    pages_missing_renderer: list = []

    doc_title: str | None = None
    sections: list = []
    landscape_pages = 0
    total_chars = 0
    page_count = 0
    with pdfplumber.open(path) as pdf:
        for index, page in enumerate(pdf.pages):
            page_no = index + 1
            page_count += 1
            if float(page.width) > float(page.height):
                landscape_pages += 1
            text = page.extract_text() or ""
            total_chars += len(text)

            if _pdf_page_is_scanned(text, page):
                if renderer is None and not (raster_map or {}).get(index):
                    pages_missing_renderer.append(page_no)
                section = _pdf_scanned_section(
                    page_no,
                    index,
                    renderer,
                    (raster_map or {}).get(index, []),
                    max_bytes,
                    cov,
                )
                if section is not None:
                    sections.append(section)
                continue

            heading = _pdf_heading(page)
            section = _pdf_page_section(text, page_no, heading_override=heading)
            page_area = float(page.width) * float(page.height) or 1.0

            table_bboxes: list = []
            try:
                found_tables = page.find_tables() or []
            except Exception:  # noqa: BLE001 - table detection is best-effort
                found_tables = []
            for table in found_tables:
                table_bboxes.append(tuple(table.bbox))
                block = _table_from_grid(table.extract())
                if block is not None:
                    cov["tables"] += 1
                    section["blocks"].append(block)

            text_lines = _pdf_text_lines(page)
            plumber_images = getattr(page, "images", None) or []
            image_bboxes = [
                (image["x0"], image["top"], image["x1"], image["bottom"])
                for image in plumber_images
            ]

            visuals: list = []  # (sort_top, block)
            page_rasters = (raster_map or {}).get(index, [])
            bboxes_by_order = (
                image_bboxes if len(image_bboxes) == len(page_rasters) else []
            )
            for raster_index, entry in enumerate(page_rasters):
                if entry["hash"] in repeated:
                    if entry["hash"] in emitted_repeats:
                        cov["images_found"] += 1
                        cov["images_skipped"] += 1
                        continue
                    emitted_repeats.add(entry["hash"])
                    pages = sorted(repeated[entry["hash"]])
                    cov["skip_reasons"].append(
                        f"repeated-asset: image on {len(pages)} pages kept once "
                        f"(page {pages[0] + 1})"
                    )
                bbox = bboxes_by_order[raster_index] if bboxes_by_order else None
                caption = _match_caption(bbox, text_lines) if bbox else None
                frac = (
                    (bbox[2] - bbox[0]) * (bbox[3] - bbox[1]) / page_area
                    if bbox
                    else None
                )
                block = _image_block(
                    entry["blob"],
                    _sniff_image_type(entry["blob"]),
                    entry["name"],
                    max_bytes,
                    cov,
                    origin="embedded-raster",
                    page=page_no,
                    caption=caption,
                    page_fraction=frac,
                )
                if block is not None:
                    visuals.append((bbox[1] if bbox else 1e9, block))

            regions = _pdf_figure_regions(page, table_bboxes, image_bboxes)
            if regions and renderer is None:
                cov["vector_regions_skipped"] += len(regions)
                pages_missing_renderer.append(page_no)
            elif regions:
                for region_no, bbox in enumerate(regions, start=1):
                    crop = _expand_region_with_labels(page, bbox)
                    png = _pdf_render_region(renderer, index, crop)
                    if png is None:
                        cov["vector_regions_skipped"] += 1
                        continue
                    caption = _match_caption(bbox, text_lines)
                    frac = (bbox[2] - bbox[0]) * (bbox[3] - bbox[1]) / page_area
                    block = _image_block(
                        png,
                        "image/png",
                        f"Figure region {region_no} on page {page_no}",
                        max_bytes,
                        cov,
                        origin="rasterized-region",
                        page=page_no,
                        caption=caption,
                        page_fraction=frac,
                    )
                    if block is not None:
                        cov["vector_regions_rasterized"] += 1
                        visuals.append((bbox[1], block))

            visuals.sort(key=lambda item: item[0])
            section["blocks"].extend(block for _top, block in visuals)

            if doc_title is None and not section["heading"].startswith("Page "):
                doc_title = section["heading"]
            sections.append(section)

    if pages_missing_renderer:
        pages = ", ".join(str(number) for number in sorted(set(pages_missing_renderer)))
        print(
            f"Warning: pypdfium2 not installed; skipped rasterizing figure "
            f"regions / scanned pages on page(s) {pages}. "
            "Please run: pip install pypdfium2",
            file=sys.stderr,
        )
    if page_count and (
        landscape_pages / page_count >= DECK_LANDSCAPE_RATIO
        and total_chars / page_count < DECK_MAX_AVG_CHARS
    ):
        cov["_deck_like"] = True
    return doc_title or Path(path).stem, sections


# --- universal ingestion: text / code / markdown / csv / image -------------


def _read_text_file(path: str, max_text_bytes: int) -> tuple[str, bool, bool]:
    """Read a text file as UTF-8 (latin-1 fallback), capped to max_text_bytes.

    Returns (text, truncated, decode_fallback).
    """
    raw = Path(path).read_bytes()
    decode_fallback = False
    try:
        text = raw.decode("utf-8")
    except UnicodeDecodeError:
        text = raw.decode("latin-1", errors="replace")
        decode_fallback = True
    truncated = False
    if len(text) > max_text_bytes:
        text = text[:max_text_bytes]
        truncated = True
    return text, truncated, decode_fallback


def _one_section(heading: str, blocks: list, kind: str = "content") -> dict:
    """A single source-index-0 section (the caller fixes source_index later)."""
    return {
        "heading": heading,
        "subheading": None,
        "kind": kind,
        "source_index": 0,
        "blocks": blocks,
    }


def _extract_code(
    path: str, max_text_bytes: int, cov: dict, rel_path: str
) -> tuple[str, list]:
    """A source-code / config file -> one section with a single code block."""
    text, truncated, fallback = _read_text_file(path, max_text_bytes)
    if fallback:
        cov["skip_reasons"].append(f"decode-fallback (latin-1): {rel_path}")
    if truncated:
        text = text.rstrip() + "\n\n... [truncated at the text-byte cap] ..."
    block = {
        "type": "code",
        "text": text,
        "language": _code_language(path),
        "path": rel_path,
    }
    if truncated:
        block["truncated"] = True
    return rel_path, [_one_section(rel_path, [block])]


def _extract_text(
    path: str, max_text_bytes: int, cov: dict, rel_path: str
) -> tuple[str, list]:
    """A plain-text file -> one section of blank-line-separated paragraphs."""
    text, truncated, fallback = _read_text_file(path, max_text_bytes)
    if fallback:
        cov["skip_reasons"].append(f"decode-fallback (latin-1): {rel_path}")
    lines = text.splitlines()
    heading = rel_path
    start = 0
    # Best-effort reStructuredText underlined title.
    if (
        len(lines) >= 2
        and lines[0].strip()
        and len(lines[1].strip()) >= 3
        and set(lines[1].strip()) <= set("=-~^\"'#*+`")
    ):
        heading = lines[0].strip()
        start = 2
    paragraphs: list = []
    current: list = []
    for line in lines[start:]:
        if line.strip():
            current.append(line.strip())
        elif current:
            paragraphs.append(" ".join(current))
            current = []
    if current:
        paragraphs.append(" ".join(current))
    blocks = [{"type": "paragraph", "text": para} for para in paragraphs if para]
    if truncated:
        blocks.append(
            {"type": "paragraph", "text": "[truncated at the text-byte cap]"}
        )
    return heading, [_one_section(heading, blocks)]


def _md_flush_paragraph(buffer: list, blocks: list) -> None:
    text = " ".join(buffer).strip()
    if text:
        blocks.append({"type": "paragraph", "text": text})
    buffer.clear()


def _md_local_image(
    line: str, base_dir: Path, max_bytes: int, cov: dict
) -> dict | None:
    """A standalone Markdown image line -> an image block or a note paragraph."""
    match = re.match(
        r"^!\[([^\]]*)\]\(([^)\s]+)(?:\s+\"[^\"]*\")?\)\s*$", line.strip()
    )
    if not match:
        return None
    alt, src = match.group(1), match.group(2)
    if src.startswith(("http://", "https://", "//")):
        return {"type": "paragraph", "text": f"[Image: {alt or src} (remote: {src})]"}
    target = (base_dir / src).resolve()
    try:
        blob = target.read_bytes()
    except OSError:
        return {"type": "paragraph", "text": f"[Image not found: {src}]"}
    if target.suffix.lower() == ".svg":
        return _svg_image_block(
            blob, alt or target.name, max_bytes, cov, origin="inline-image"
        )
    return _image_block(
        blob,
        _sniff_image_type(blob),
        alt or target.name,
        max_bytes,
        cov,
        origin="inline-image",
    )


def _extract_markdown(
    path: str, max_text_bytes: int, max_bytes: int, cov: dict, rel_path: str
) -> tuple[str, list]:
    """Parse Markdown into heading-delimited sections (intentionally minimal)."""
    text, truncated, fallback = _read_text_file(path, max_text_bytes)
    if fallback:
        cov["skip_reasons"].append(f"decode-fallback (latin-1): {rel_path}")
    base_dir = Path(path).parent
    lines = text.splitlines()
    doc_title: str | None = None
    sections: list = []
    current = _one_section(rel_path, [])
    para: list = []
    pending_bullets: list = []

    def close_para() -> None:
        _md_flush_paragraph(para, current["blocks"])

    def close_bullets() -> None:
        if pending_bullets:
            current["blocks"].append(
                {"type": "bullets", "items": list(pending_bullets)}
            )
            pending_bullets.clear()

    def start_section(heading: str) -> None:
        nonlocal current
        close_para()
        close_bullets()
        if current["blocks"] or current["heading"] != rel_path:
            sections.append(current)
        current = _one_section(heading, [])

    index = 0
    while index < len(lines):
        line = lines[index]
        stripped = line.strip()
        fence = re.match(r"^(`{3,}|~{3,})\s*([\w+#-]*)\s*$", stripped)
        if fence:
            close_para()
            close_bullets()
            marker = fence.group(1)[0] * 3
            body: list = []
            index += 1
            while index < len(lines) and not lines[index].strip().startswith(marker):
                body.append(lines[index])
                index += 1
            index += 1  # skip the closing fence
            current["blocks"].append(
                {"type": "code", "text": "\n".join(body), "language": fence.group(2)}
            )
            continue
        atx = re.match(r"^(#{1,6})\s+(.*?)\s*#*\s*$", stripped)
        if atx:
            heading = atx.group(2).strip()
            if doc_title is None:
                doc_title = heading
            start_section(heading)
            index += 1
            continue
        if (
            stripped
            and index + 1 < len(lines)
            and set(lines[index + 1].strip()) in ({"="}, {"-"})
            and len(lines[index + 1].strip()) >= 3
            and stripped[0] not in "-*+|>#"
        ):
            if doc_title is None:
                doc_title = stripped
            start_section(stripped)
            index += 2
            continue
        image_block = _md_local_image(stripped, base_dir, max_bytes, cov)
        if image_block is not None:
            close_para()
            close_bullets()
            current["blocks"].append(image_block)
            index += 1
            continue
        if (
            "|" in stripped
            and index + 1 < len(lines)
            and "-" in lines[index + 1]
            and re.match(r"^\s*\|?[\s:|-]+\|?\s*$", lines[index + 1])
        ):
            close_para()
            close_bullets()
            table_lines = [stripped]
            index += 2  # skip header + separator
            while index < len(lines) and "|" in lines[index] and lines[index].strip():
                table_lines.append(lines[index].strip())
                index += 1
            grid = [
                [cell.strip() for cell in row.strip().strip("|").split("|")]
                for row in table_lines
            ]
            block = _table_from_grid(grid)
            if block is not None:
                cov["tables"] += 1
                current["blocks"].append(block)
            continue
        bullet = re.match(r"^(\s*)(?:[-*+]|\d+\.)\s+(.*)$", line)
        if bullet:
            close_para()
            pending_bullets.append(
                {"text": bullet.group(2).strip(), "depth": len(bullet.group(1)) // 2}
            )
            index += 1
            continue
        if not stripped:
            close_para()
            close_bullets()
        else:
            close_bullets()
            para.append(stripped)
        index += 1

    close_para()
    close_bullets()
    if current["blocks"] or current["heading"] != rel_path or not sections:
        sections.append(current)
    if truncated and sections:
        sections[-1]["blocks"].append(
            {"type": "paragraph", "text": "[truncated at the text-byte cap]"}
        )
    return doc_title or rel_path, sections


def _coerce_cell(value: str) -> object:
    """Turn a numeric-looking CSV cell into a float; keep everything else."""
    stripped = value.strip().replace(",", "")
    if _NUMBER_RE.match(stripped):
        try:
            return float(stripped)
        except ValueError:
            return value
    return value


def _extract_csv(
    path: str, max_text_bytes: int, cov: dict, rel_path: str
) -> tuple[str, list]:
    """A CSV / TSV file -> a chart (numeric) or table block via the grid logic."""
    import csv
    import io

    text, truncated, fallback = _read_text_file(path, max_text_bytes)
    if fallback:
        cov["skip_reasons"].append(f"decode-fallback (latin-1): {rel_path}")
    if Path(path).suffix.lower() == ".tsv":
        delimiter = "\t"
    else:
        try:
            delimiter = csv.Sniffer().sniff(text[:4096], delimiters=",\t;|").delimiter
        except csv.Error:
            delimiter = ","
    rows: list = []
    capped = truncated
    for row in csv.reader(io.StringIO(text), delimiter=delimiter):
        if len(rows) >= CSV_MAX_ROWS:
            capped = True
            break
        rows.append([_coerce_cell(cell) for cell in row[:CSV_MAX_COLS]])
    if capped:
        cov["skip_reasons"].append(f"csv-cap: {rel_path} truncated for size")
    block = _grid_to_block(rows) if rows else None
    if block is not None and block["type"] == "table":
        cov["tables"] += 1
    return Path(path).stem, [
        _one_section(rel_path, [block] if block is not None else [], kind="data")
    ]


def _svg_image_block(
    data: bytes, alt: str, max_bytes: int, cov: dict, origin: str = "standalone-image"
) -> dict | None:
    """Embed SVG markup as an image/svg+xml data URI (never rasterized)."""
    cov["images_found"] += 1
    if len(data) > max_bytes:
        cov["images_skipped"] += 1
        cov["skip_reasons"].append(f"over-budget svg: {alt}")
        return None
    encoded = base64.b64encode(data).decode("ascii")
    cov["images_kept"] += 1
    return {
        "type": "image",
        "data_uri": f"data:image/svg+xml;base64,{encoded}",
        "alt": alt or "Image",
        "origin": origin,
    }


def _extract_image_file(
    path: str, max_bytes: int, cov: dict, rel_path: str
) -> tuple[str, list]:
    """A standalone image file -> one image section (origin standalone-image)."""
    data = Path(path).read_bytes()
    name = Path(path).stem
    if Path(path).suffix.lower() == ".svg":
        block = _svg_image_block(data, name, max_bytes, cov)
    else:
        block = _image_block(
            data,
            _sniff_image_type(data),
            name,
            max_bytes,
            cov,
            origin="standalone-image",
        )
    return name, [_one_section(rel_path, [block] if block is not None else [], kind="image")]


# --- directory / repository walk --------------------------------------------


def _note(report: dict, message: str, limit: int = 40) -> None:
    """Append a bounded human-readable note to a walk report."""
    if len(report["notes"]) < limit:
        report["notes"].append(message)


def _is_binary(path: str) -> bool:
    """True when a file looks binary (a NUL byte or many non-text bytes)."""
    try:
        with open(path, "rb") as handle:
            chunk = handle.read(4096)
    except OSError:
        return True
    if not chunk:
        return False
    if b"\x00" in chunk:
        return True
    printable = set(range(32, 127)) | set(b"\n\r\t\f\b")
    nontext = sum(1 for byte in chunk if byte not in printable)
    return nontext / len(chunk) > 0.30


def _load_gitignore(root: Path) -> list:
    """Best-effort .gitignore patterns (limits documented in the runbook)."""
    patterns: list = []
    gitignore = root / ".gitignore"
    if gitignore.is_file():
        try:
            content = gitignore.read_text(encoding="utf-8", errors="replace")
        except OSError:
            return patterns
        for line in content.splitlines():
            line = line.strip()
            if line and not line.startswith("#") and not line.startswith("!"):
                patterns.append(line.rstrip("/"))
    return patterns


def _gitignored(rel_posix: str, name: str, patterns: list) -> bool:
    """Approximate .gitignore match (no negation, no full '**' semantics)."""
    import fnmatch

    for pattern in patterns:
        if pattern.startswith("/"):
            anchored = pattern.lstrip("/")
            if (
                rel_posix == anchored
                or rel_posix.startswith(anchored + "/")
                or fnmatch.fnmatch(rel_posix, anchored)
            ):
                return True
        elif (
            fnmatch.fnmatch(name, pattern)
            or fnmatch.fnmatch(rel_posix, pattern)
            or fnmatch.fnmatch(rel_posix, "*/" + pattern)
            or ("/" + pattern + "/") in ("/" + rel_posix + "/")
        ):
            return True
    return False


def _walk_directory(root: str, max_files: int) -> tuple[list, dict]:
    """Recursively collect supported files under a directory (deterministic).

    Ignored dirs, lockfiles, .gitignore matches, and binaries are excluded;
    beyond the file-count cap, extra files are dropped. Returns
    (sorted file paths, walk-report dict).
    """
    root_path = Path(root)
    patterns = _load_gitignore(root_path)
    report = {
        "root": root_path.name,
        "files_included": 0,
        "gitignored": 0,
        "binary_skipped": 0,
        "dirs_ignored": 0,
        "file_count_capped": 0,
        "notes": [],
    }
    included: list = []
    for dirpath, dirnames, filenames in os.walk(root):
        kept: list = []
        for name in sorted(dirnames):
            if name.lower() in IGNORE_DIRS:
                report["dirs_ignored"] += 1
            else:
                kept.append(name)
        dirnames[:] = kept
        for filename in sorted(filenames):
            full = Path(dirpath) / filename
            rel = full.relative_to(root_path).as_posix()
            if filename.lower() in IGNORE_BASENAMES:
                _note(report, f"ignored lockfile: {rel}")
                continue
            if _gitignored(rel, filename, patterns):
                report["gitignored"] += 1
                continue
            fmt = _detect_format(str(full))
            if fmt is None:
                continue
            if fmt in TEXT_FORMATS and _is_binary(str(full)):
                report["binary_skipped"] += 1
                _note(report, f"binary-skipped: {rel}")
                continue
            included.append((rel, str(full)))
    included.sort(key=lambda item: item[0])
    if len(included) > max_files:
        dropped = len(included) - max_files
        report["file_count_capped"] = dropped
        _note(report, f"file-count-cap: kept {max_files}, skipped {dropped}")
        included = included[:max_files]
    report["files_included"] = len(included)
    return [full for _rel, full in included], report


def _build_tree(rel_paths: list, root_name: str) -> dict:
    """Nested {name, kind, children} tree of the ingested files (sorted)."""
    root = {"name": root_name, "kind": "dir", "children": []}
    for rel in rel_paths:
        parts = rel.split("/")
        node = root
        for depth, part in enumerate(parts):
            is_file = depth == len(parts) - 1
            child = next(
                (
                    c
                    for c in node["children"]
                    if c["name"] == part and (c["kind"] == "file") == is_file
                ),
                None,
            )
            if child is None:
                child = {"name": part, "kind": "file" if is_file else "dir"}
                if not is_file:
                    child["children"] = []
                node["children"].append(child)
            if not is_file:
                node = child
    _sort_tree(root)
    return root


def _sort_tree(node: dict) -> None:
    children = node.get("children")
    if not children:
        return
    children.sort(key=lambda c: (c["kind"] != "dir", c["name"].lower()))
    for child in children:
        _sort_tree(child)


# --- dispatch and merge ----------------------------------------------------


def _detect_format(path: str) -> str | None:
    name = Path(path).name.lower()
    if name in SPECIAL_BASENAMES:
        return SPECIAL_BASENAMES[name][0]
    return EXTENSION_FORMATS.get(Path(path).suffix.lower())


def _code_language(path: str) -> str:
    """The syntax-highlight language for a source-code / config file."""
    name = Path(path).name.lower()
    if name in SPECIAL_BASENAMES:
        return SPECIAL_BASENAMES[name][1]
    return CODE_LANGUAGES.get(Path(path).suffix.lower(), "")


def _expand_inputs(
    paths: list, max_files: int
) -> tuple[list, dict | None, str | None]:
    """Expand inputs to a file list, walking any directory recursively.

    Returns (files, walk_report, repo_root). repo_root is set (enabling
    repository assembly) only when the sole input is a single directory.
    """
    single_dir = len(paths) == 1 and Path(paths[0]).is_dir()
    files: list = []
    walk_report: dict | None = None
    repo_root: str | None = None
    for raw in paths:
        path = Path(raw)
        if path.is_dir():
            walked, report = _walk_directory(str(path), max_files)
            files.extend(walked)
            if walk_report is None:
                walk_report = report
            if single_dir:
                repo_root = str(path)
        else:
            files.append(str(path))
    return files, walk_report, repo_root


def _extract_one(
    path: str,
    fmt: str,
    max_bytes: int,
    max_text_bytes: int,
    cov: dict,
    rel_path: str,
) -> tuple[str, list]:
    if fmt == "pptx":
        return _extract_pptx(path, max_bytes, cov)
    if fmt == "docx":
        return _extract_docx(path, max_bytes, cov)
    if fmt == "xlsx":
        return _extract_xlsx(path, cov)
    if fmt == "pdf":
        return _extract_pdf(path, max_bytes, cov)
    if fmt == "code":
        return _extract_code(path, max_text_bytes, cov, rel_path)
    if fmt == "markdown":
        return _extract_markdown(path, max_text_bytes, max_bytes, cov, rel_path)
    if fmt == "text":
        return _extract_text(path, max_text_bytes, cov, rel_path)
    if fmt == "csv":
        return _extract_csv(path, max_text_bytes, cov, rel_path)
    if fmt == "image":
        return _extract_image_file(path, max_bytes, cov, rel_path)
    raise ValueError(f"unsupported format: {fmt}")


def _agenda_section(headings: list) -> dict:
    return {
        "heading": "Agenda",
        "subheading": None,
        "kind": "section-break",
        "source_index": 0,
        "blocks": [
            {
                "type": "bullets",
                "items": [{"text": head, "depth": 0} for head in headings],
            }
        ],
    }


def _title_section(title: str) -> dict:
    return {
        "heading": title,
        "subheading": None,
        "kind": "title",
        "source_index": 0,
        "blocks": [],
    }


def _repo_category(fmt: str, rel_path: str) -> int:
    """Ordering rank for repository assembly (README, docs, code, data, ...)."""
    base = rel_path.rsplit("/", 1)[-1].lower()
    if fmt in ("markdown", "text"):
        return 0 if base.startswith("readme") else 1
    if fmt == "code":
        return 2
    if fmt == "csv":
        return 3
    if fmt == "image":
        return 4
    return 5


def _top_level_dir(rel_path: str) -> str:
    return rel_path.split("/", 1)[0] if "/" in rel_path else "(root)"


def _assemble_single(per_source: list, title_override: str | None) -> tuple[list, str]:
    _index, fmt, title, sections, _rel = per_source[0]
    model_title = title_override or title
    out: list = []
    if fmt in ("docx", "pdf", "xlsx"):
        out.append(_title_section(model_title))
        headings = [
            section["heading"]
            for section in sections
            if section["kind"] in ("content", "data") and section["heading"]
        ]
        if fmt in ("docx", "pdf") and len(headings) >= 2:
            out.append(_agenda_section(headings))
    for section in sections:
        section["source_index"] = 0
        out.append(section)
    return out, model_title


def _assemble_multi(
    per_source: list, sources: list, title_override: str | None
) -> tuple[list, str]:
    umbrella = title_override or sources[0]["title"]
    out: list = [
        {
            "heading": umbrella,
            "subheading": "Combined sources",
            "kind": "title",
            "source_index": 0,
            "blocks": [
                {
                    "type": "bullets",
                    "items": [{"text": src["title"], "depth": 0} for src in sources],
                }
            ],
        }
    ]
    for index, _fmt, title, sections, _rel in per_source:
        out.append(
            {
                "heading": title,
                "subheading": None,
                "kind": "section-break",
                "source_index": index,
                "blocks": [],
            }
        )
        for section in sections:
            section["source_index"] = index
            out.append(section)
    return out, umbrella


def _repo_readme_blurb(per_source: list) -> str | None:
    """The first non-empty paragraph of a README, for the overview section."""
    for _index, fmt, _title, sections, rel in per_source:
        base = rel.rsplit("/", 1)[-1].lower()
        if fmt in ("markdown", "text") and base.startswith("readme"):
            for section in sections:
                for block in section["blocks"]:
                    if block["type"] == "paragraph" and block["text"].strip():
                        return block["text"].strip()
    return None


def _assemble_repository(
    per_source: list, root_path: Path, title_override: str | None
) -> tuple[list, str]:
    model_title = title_override or root_path.name
    top_entries: set = set()
    for _index, _fmt, _title, _sections, rel in per_source:
        tld = _top_level_dir(rel)
        top_entries.add((tld + "/") if tld != "(root)" else rel)
    overview_blocks: list = []
    blurb = _repo_readme_blurb(per_source)
    if blurb:
        overview_blocks.append({"type": "paragraph", "text": blurb})
    if top_entries:
        overview_blocks.append(
            {
                "type": "bullets",
                "items": [
                    {"text": entry, "depth": 0} for entry in sorted(top_entries)[:24]
                ],
            }
        )
    out: list = [
        {
            "heading": model_title,
            "subheading": "Project overview",
            "kind": "overview",
            "source_index": 0,
            "blocks": overview_blocks,
        }
    ]
    ordered = sorted(
        per_source,
        key=lambda entry: (
            _repo_category(entry[1], entry[4]),
            _top_level_dir(entry[4]),
            entry[4],
        ),
    )
    last_code_dir: str | None = None
    for index, fmt, _title, sections, rel in ordered:
        if fmt == "code":
            tld = _top_level_dir(rel)
            if tld != last_code_dir:
                out.append(
                    {
                        "heading": tld if tld != "(root)" else "Top-level files",
                        "subheading": None,
                        "kind": "section-break",
                        "source_index": 0,
                        "blocks": [],
                    }
                )
                last_code_dir = tld
        for section in sections:
            section["source_index"] = index
            out.append(section)
    return out, model_title


def build_model(
    paths: list,
    max_bytes: int = DEFAULT_MAX_IMAGE_BYTES,
    title_override: str | None = None,
    max_text_bytes: int = DEFAULT_MAX_TEXT_BYTES,
    max_files: int = DEFAULT_MAX_FILES,
) -> dict:
    """Extract every input into one merged, deterministic content model."""
    files, walk_report, repo_root = _expand_inputs(paths, max_files)
    root_path = Path(repo_root) if repo_root else None

    sources: list = []
    coverage: list = []
    per_source: list = []  # (index, fmt, title, sections, rel_path)
    for path in files:
        fmt = _detect_format(path)
        if fmt is None:
            print(f"Warning: skipping unsupported file: {path}", file=sys.stderr)
            continue
        if not Path(path).is_file():
            print(f"Warning: input not found: {path}", file=sys.stderr)
            continue
        if root_path is not None:
            try:
                rel_path = Path(path).relative_to(root_path).as_posix()
            except ValueError:
                rel_path = Path(path).name
        else:
            rel_path = Path(path).name
        cov = _new_coverage(rel_path)
        try:
            detected_title, sections = _extract_one(
                path, fmt, max_bytes, max_text_bytes, cov, rel_path
            )
        except SystemExit:
            raise  # missing-dependency exit must propagate
        except Exception as exc:  # noqa: BLE001 - logged, then this file is skipped
            print(f"Warning: failed to extract {path}: {exc}", file=sys.stderr)
            continue
        index = len(sources)
        title = detected_title or Path(path).stem
        source = {"path": rel_path, "format": fmt, "title": title}
        if cov.pop("_deck_like", False):
            source["deck_like"] = True
        sources.append(source)
        coverage.append(cov)
        per_source.append((index, fmt, title, sections, rel_path))

    if not sources:
        print(
            "Error: no supported input found (.pptx/.docx/.xlsx/.pdf, "
            "code/markdown/text/csv, or images).",
            file=sys.stderr,
        )
        raise SystemExit(2)

    tree: dict | None = None
    if root_path is not None:
        out_sections, model_title = _assemble_repository(
            per_source, root_path, title_override
        )
        tree = _build_tree([entry[4] for entry in per_source], root_path.name)
    elif len(sources) > 1:
        out_sections, model_title = _assemble_multi(per_source, sources, title_override)
    else:
        out_sections, model_title = _assemble_single(per_source, title_override)

    coverage_obj: dict = {"per_source": coverage}
    if walk_report is not None:
        coverage_obj["walk"] = walk_report
    model: dict = {
        "schema_version": 2,
        "title": model_title,
        "sources": sources,
        "sections": out_sections,
        "coverage": coverage_obj,
    }
    placeholders = _scan_placeholders(out_sections)
    if placeholders:
        model["placeholders"] = placeholders
    if tree is not None:
        model["tree"] = tree
    return model


# Unfilled author placeholders ("enrolled xx subjects", "[insert Q3 figure]")
# ship inside real documents more often than one would hope. Detect them so the
# agent can ASK the user for the real value during the design intake instead of
# discovering the gap after delivery - and never silently invent the number.
_PLACEHOLDER_RES = [
    re.compile(r"(?<![\w.])[xX]{2,}(?![\w.])"),
    re.compile(r"\[(?:insert|tbd|placeholder|todo|fill[ -]?in)[^\]]{0,60}\]", re.IGNORECASE),
    re.compile(r"<(?:insert|tbd|placeholder|todo)[^>]{0,60}>", re.IGNORECASE),
]
_PLACEHOLDER_BLOCK_TYPES = {"paragraph", "heading", "bullets", "table"}
_PLACEHOLDER_CAP = 50


def _scan_placeholders(sections: list) -> list:
    """Scan document-text blocks for unfilled author placeholders.

    Deliberately conservative (repeated x tokens and bracketed insert/TBD
    markers only) and scoped to prose block types, so code listings from a
    repository walk cannot flood the report with hex strings or templates.
    """
    found: list = []
    for s_idx, section in enumerate(sections):
        for block in section.get("blocks", []):
            if block.get("type") not in _PLACEHOLDER_BLOCK_TYPES:
                continue
            texts: list = []
            if isinstance(block.get("text"), str):
                texts.append(block["text"])
            for item in block.get("items", []) or []:
                if isinstance(item, dict) and isinstance(item.get("text"), str):
                    texts.append(item["text"])
                elif isinstance(item, str):
                    texts.append(item)
            for row in block.get("rows", []) or []:
                for cell in row or []:
                    if isinstance(cell, str):
                        texts.append(cell)
            for text in texts:
                for rx in _PLACEHOLDER_RES:
                    for match in rx.finditer(text):
                        start = max(0, match.start() - 40)
                        end = min(len(text), match.end() + 40)
                        found.append(
                            {
                                "section_index": s_idx,
                                "section_title": section.get("title") or "",
                                "block_type": block.get("type"),
                                "token": match.group(0),
                                "context": text[start:end].strip(),
                            }
                        )
                        if len(found) >= _PLACEHOLDER_CAP:
                            return found
    return found


def main(argv: list | None = None) -> int:
    parser = argparse.ArgumentParser(
        description=(
            "Extract one or more inputs into the normalized content-model JSON "
            "consumed by build_presentation.py. Inputs may be documents "
            "(.pptx/.docx/.xlsx/.pdf), source code / config, Markdown / text, "
            "CSV / TSV, standalone images, or a directory / repository (walked "
            "recursively). Local-only parsing; no network calls."
        )
    )
    parser.add_argument(
        "inputs", nargs="+", help="Input file(s), folder(s), or a repository."
    )
    parser.add_argument(
        "-o", "--out", required=True, help="Output content-model JSON path."
    )
    parser.add_argument(
        "--max-image-bytes",
        type=int,
        default=DEFAULT_MAX_IMAGE_BYTES,
        help="Per-image byte budget before downscale/skip (default 2000000).",
    )
    parser.add_argument(
        "--max-text-bytes",
        type=int,
        default=DEFAULT_MAX_TEXT_BYTES,
        help="Per-text-file byte cap before truncation (default 300000).",
    )
    parser.add_argument(
        "--max-files",
        type=int,
        default=DEFAULT_MAX_FILES,
        help="Max files ingested from a directory walk (default 400).",
    )
    parser.add_argument(
        "--title", default=None, help="Override the synthesized presentation title."
    )
    args = parser.parse_args(argv)

    model = build_model(
        args.inputs,
        args.max_image_bytes,
        args.title,
        args.max_text_bytes,
        args.max_files,
    )
    out_path = Path(args.out)
    if out_path.parent and not out_path.parent.exists():
        out_path.parent.mkdir(parents=True, exist_ok=True)
    with out_path.open("w", encoding="utf-8") as handle:
        json.dump(model, handle, ensure_ascii=False, indent=2)
        handle.write("\n")
    walk = model.get("coverage", {}).get("walk")
    if walk is not None:
        print(
            f"Walked {walk['root']}: {walk['files_included']} file(s) included, "
            f"{walk['gitignored']} gitignored, {walk['binary_skipped']} binary, "
            f"{walk['dirs_ignored']} dir(s) ignored, "
            f"{walk['file_count_capped']} over the file cap.",
            file=sys.stderr,
        )
    placeholders = model.get("placeholders") or []
    if placeholders:
        print(
            f"NOTE: {len(placeholders)} unfilled source placeholder(s) detected "
            f"(see the model's 'placeholders' list) - surface them to the user "
            f"during the design intake; never invent the missing value.",
            file=sys.stderr,
        )
        for entry in placeholders[:5]:
            print(
                f"  - [{entry['token']}] in section {entry['section_index']}: "
                f"...{entry['context'][:90]}...",
                file=sys.stderr,
            )
    print(
        f"Wrote {out_path} ({len(model['sources'])} source(s), "
        f"{len(model['sections'])} section(s)).",
        file=sys.stderr,
    )
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
